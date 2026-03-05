#!/usr/bin/env python3
"""End-to-end test for slides build-rich: creates a sample .pptx locally."""
import json
import os
import sys
import tempfile

# Add parent dir so we can import the module
sys.path.insert(0, os.path.dirname(__file__))

from google_workspace import _build_rich_presentation

SAMPLE_SPEC = {
    "title": "Test Presentation",
    "theme": {
        "dark": "1A1A2E",
        "accent": "E94560",
        "accent2": "0F3460",
        "lightBg": "F5F5F5",
        "medGray": "666666",
    },
    "slides": [
        {
            "type": "title_slide",
            "title": "Q1 Performance Review",
            "subtitle": "March 2026 - Test Brand",
            "notes": "Welcome everyone to the Q1 review",
        },
        {
            "type": "section",
            "title": "Executive Summary",
        },
        {
            "type": "kpi_row",
            "title": "Key Metrics",
            "metrics": [
                {"value": "$42K", "label": "Revenue", "bg": "0F3460"},
                {"value": "3.8x", "label": "ROAS", "bg": "2E7D32"},
                {"value": "186", "label": "Orders", "bg": "1565C0"},
                {"value": "+10.5%", "label": "Growth", "bg": "E94560"},
            ],
        },
        {
            "type": "content",
            "title": "Key Findings",
            "body": (
                "Revenue grew 10.5% this quarter driven by strong organic performance "
                "and improved paid media efficiency. The top 5 ads accounted for 60% of "
                "total revenue, with UGC content outperforming studio creative by 1.8x."
            ),
        },
        {
            "type": "bullets",
            "title": "Growth Drivers",
            "items": [
                "Organic traffic up 15% month-over-month",
                "Paid ROAS improved from 3.4x to 3.8x",
                "Email conversion rate steady at 4.2%",
                "New customer acquisition cost down 8%",
            ],
        },
        {
            "type": "table",
            "title": "Ad Performance Summary",
            "headers": ["Ad Name", "ROAS", "Spend", "Revenue", "CTR"],
            "rows": [
                ["Morning Routine UGC", "4.2x", "$1,200", "$5,040", "2.8%"],
                ["Product Hero Shot", "2.8x", "$800", "$2,240", "1.9%"],
                ["Before/After Carousel", "3.5x", "$950", "$3,325", "2.3%"],
                ["Lifestyle Flat Lay", "2.1x", "$600", "$1,260", "1.5%"],
            ],
            "style": {"headerBg": "0F3460", "zebra": True},
        },
        {
            "type": "chart",
            "title": "Revenue Trend",
            "chartType": "bar",
            "categories": ["January", "February", "March"],
            "series": [
                {"name": "Revenue", "values": [32000, 38000, 42000]},
                {"name": "Target", "values": [35000, 37000, 40000]},
            ],
        },
        {
            "type": "two_column",
            "title": "Strategy vs Results",
            "left": {
                "heading": "Q1 Strategy",
                "body": (
                    "Focus on scaling UGC content creators and testing new ad formats. "
                    "Reduce dependency on single hero creatives."
                ),
            },
            "right": {
                "heading": "Q1 Results",
                "body": (
                    "UGC drove 3.2x ROAS vs 2.1x for studio. Successfully diversified "
                    "creative mix with 12 new ad variants."
                ),
            },
        },
        {
            "type": "blank",
            "background": "FFFFFF",
            "elements": [
                {
                    "type": "text",
                    "text": "Custom Layout Example",
                    "x": 0.8,
                    "y": 0.5,
                    "width": 11.7,
                    "height": 0.8,
                    "fontSize": 28,
                    "bold": True,
                    "color": "1A1A2E",
                },
                {
                    "type": "text",
                    "text": "This slide uses the blank type with manually positioned elements.",
                    "x": 0.8,
                    "y": 1.5,
                    "width": 11.7,
                    "height": 0.6,
                    "fontSize": 16,
                    "color": "666666",
                },
                {
                    "type": "shape",
                    "x": 0.8,
                    "y": 2.3,
                    "width": 11.7,
                    "height": 0.06,
                    "color": "E94560",
                },
            ],
        },
        {
            "type": "closing",
            "title": "Thank You",
            "subtitle": "Questions & Discussion",
        },
    ],
}


def main():
    output_path = os.path.join(tempfile.mkdtemp(), "test_presentation.pptx")

    print(f"Building test presentation...")
    _build_rich_presentation(SAMPLE_SPEC, output_path)

    size = os.path.getsize(output_path)
    print(f"Created: {output_path}")
    print(f"Size: {size:,} bytes")
    print(f"Slides: {len(SAMPLE_SPEC['slides'])}")

    # Verify the file is a valid pptx by opening it
    from pptx import Presentation
    prs = Presentation(output_path)
    actual_slides = len(prs.slides)
    print(f"Verified: {actual_slides} slides in output file")

    assert actual_slides == len(SAMPLE_SPEC["slides"]), (
        f"Expected {len(SAMPLE_SPEC['slides'])} slides, got {actual_slides}"
    )

    # Check each slide has content
    for i, slide in enumerate(prs.slides):
        shape_count = len(slide.shapes)
        assert shape_count > 0, f"Slide {i} has no shapes"
        print(f"  Slide {i+1}: {shape_count} shapes ({SAMPLE_SPEC['slides'][i]['type']})")

    print("\nAll tests passed.")
    print(json.dumps({"status": "ok", "path": output_path, "slides": actual_slides}))


if __name__ == "__main__":
    main()
