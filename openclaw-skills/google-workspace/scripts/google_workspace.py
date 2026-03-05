#!/usr/bin/env python3
# /// script
# requires-python = ">=3.10"
# dependencies = [
#     "google-api-python-client>=2.100.0",
#     "google-auth>=2.20.0",
#     "google-auth-oauthlib>=1.0.0",
#     "google-auth-httplib2>=0.2.0",
#     "matplotlib>=3.8.0",
#     "openpyxl>=3.1.0",
#     "pandas>=2.0.0",
#     "python-docx>=1.0.0",
#     "python-pptx>=0.6.21",
# ]
# ///
"""Google Workspace API helper — Slides, Docs, Sheets, Drive.

Subcommand-based CLI for creating polished, presentation-ready content
via Google Workspace APIs with per-profile OAuth credentials.
"""

import argparse
import io
import json
import os
import pathlib
import re
import signal
import sys
import tempfile
import uuid

# --- SIGPIPE + PipeSafe (standard openCLAW boilerplate) ---
try:
    signal.signal(signal.SIGPIPE, signal.SIG_IGN)
except (AttributeError, ValueError):
    pass

class _PipeSafe:
    def __init__(self, stream):
        self._stream = stream
    def write(self, data):
        try:
            self._stream.write(data)
        except (BrokenPipeError, OSError):
            pass
    def flush(self):
        try:
            self._stream.flush()
        except (BrokenPipeError, OSError):
            pass
    def __getattr__(self, name):
        return getattr(self._stream, name)

sys.stdout = _PipeSafe(sys.stdout)
sys.stderr = _PipeSafe(sys.stderr)

# --- Path derivation (profile-agnostic) ---
# IMPORTANT: Do NOT use .resolve() here. The script is symlinked from each
# profile directory (e.g. ~/.openclaw/skills/google-workspace/scripts/ -> repo).
# Using .resolve() would follow the symlink to the repo, breaking profile
# isolation for credentials, .env, and .oauth-client.json.
_SCRIPT_PATH = pathlib.Path(__file__).parent  # keeps symlink path
SKILL_DIR = _SCRIPT_PATH.parent
SCRIPTS_DIR = SKILL_DIR / "scripts"


def _workspace_root():
    """Derive workspace root: <root>/skills/<name>/scripts/<this> -> climb 4 dirs."""
    return _SCRIPT_PATH.parent.parent.parent


def _load_env_file():
    """Load .env files as fallback for missing env vars."""
    for env_path in [SKILL_DIR / ".env", _workspace_root() / ".env"]:
        if not env_path.exists():
            continue
        for line in env_path.read_text().splitlines():
            line = line.strip()
            if not line or line.startswith("#"):
                continue
            if "=" in line:
                key, _, value = line.partition("=")
                key = key.strip()
                value = value.strip().strip('"').strip("'")
                if key and not os.environ.get(key):
                    os.environ[key] = value


_load_env_file()

# --- Credential management ---
CREDS_DIR_NAME = "credentials"
CREDS_FILE_NAME = "google-workspace-oauth.json"
OAUTH_CLIENT_FILE = SKILL_DIR / ".oauth-client.json"

SCOPES = [
    "https://www.googleapis.com/auth/presentations",
    "https://www.googleapis.com/auth/documents",
    "https://www.googleapis.com/auth/spreadsheets",
    "https://www.googleapis.com/auth/drive",
]


def _creds_path():
    return _workspace_root() / CREDS_DIR_NAME / CREDS_FILE_NAME


def _get_credentials():
    """Load and refresh Google OAuth credentials for this profile."""
    import google.auth.transport.requests
    import google.oauth2.credentials

    cp = _creds_path()
    if not cp.exists():
        print(f"ERROR: Google Workspace not configured for this profile.", file=sys.stderr)
        print(f"Credentials expected at: {cp}", file=sys.stderr)
        print(f"Run: uv run {SCRIPTS_DIR / 'setup_oauth.py'}", file=sys.stderr)
        sys.exit(1)

    creds_data = json.loads(cp.read_text())
    creds = google.oauth2.credentials.Credentials.from_authorized_user_info(
        creds_data, scopes=SCOPES
    )

    if creds.expired and creds.refresh_token:
        creds.refresh(google.auth.transport.requests.Request())
        # Persist refreshed token
        updated = {
            "type": "authorized_user",
            "client_id": creds.client_id,
            "client_secret": creds.client_secret,
            "refresh_token": creds.refresh_token,
            "scopes": SCOPES,
            "account": creds_data.get("account", ""),
            "created_at": creds_data.get("created_at", ""),
        }
        cp.write_text(json.dumps(updated, indent=2))
        cp.chmod(0o600)

    return creds


def _build_service(api, version):
    """Build a Google API service client."""
    from googleapiclient.discovery import build
    return build(api, version, credentials=_get_credentials())


def _slides_service():
    return _build_service("slides", "v1")


def _docs_service():
    return _build_service("docs", "v1")


def _sheets_service():
    return _build_service("sheets", "v4")


def _drive_service():
    return _build_service("drive", "v3")


def _auto_share_org(file_id):
    """Share a file with the org domain (anyone at domain can view).

    Derives domain from GOG_AUTH_EMAIL env var. Silently skips if
    no email is set or if the account is a personal Gmail.
    """
    email = os.environ.get("GOG_AUTH_EMAIL", "")
    if not email or "@" not in email:
        return
    domain = email.split("@", 1)[1].lower()
    # Skip personal Gmail accounts — domain sharing doesn't work for them
    if domain in ("gmail.com", "googlemail.com"):
        return
    try:
        drive = _drive_service()
        drive.permissions().create(
            fileId=file_id,
            body={"type": "domain", "domain": domain, "role": "reader"},
            fields="id",
        ).execute()
    except Exception:
        pass  # Silent fail for personal accounts or permission errors


# ============================================================
# Utility helpers
# ============================================================

def _hex_to_rgb(hex_color):
    """Convert '#RRGGBB' to Google API RGB dict (0.0-1.0)."""
    h = hex_color.lstrip("#")
    if len(h) != 6:
        print(f"ERROR: Invalid hex color '{hex_color}'. Use '#RRGGBB'.", file=sys.stderr)
        sys.exit(1)
    r, g, b = int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16)
    return {"red": r / 255.0, "green": g / 255.0, "blue": b / 255.0}


def _rgb_color(hex_color):
    """Return an OpaqueColor dict from hex."""
    return {"opaqueColor": {"rgbColor": _hex_to_rgb(hex_color)}}


def _solid_fill(hex_color, alpha=1.0):
    """Return a solidFill dict from hex."""
    return {"solidFill": {"color": {"rgbColor": _hex_to_rgb(hex_color)}, "alpha": alpha}}


def _pt(magnitude):
    """Dimension in points."""
    return {"magnitude": magnitude, "unit": "PT"}


def _emu(magnitude):
    """Dimension in EMU."""
    return {"magnitude": magnitude, "unit": "EMU"}


def _emu_from_inches(inches):
    return int(inches * 914400)


def _emu_from_pt(pt):
    return int(pt * 12700)


def _gen_id():
    """Generate a short unique object ID for Google API elements."""
    return "g_" + uuid.uuid4().hex[:12]


def _parse_json_arg(value, name="data"):
    """Parse a JSON string argument, exit on error."""
    try:
        return json.loads(value)
    except json.JSONDecodeError as e:
        print(f"ERROR: Invalid JSON for --{name}: {e}", file=sys.stderr)
        sys.exit(1)


def _load_brand_colors():
    """Load brand colors from workspace/brand/colors.md if available."""
    colors_file = _workspace_root() / "workspace" / "brand" / "colors.md"
    if not colors_file.exists():
        return {}
    text = colors_file.read_text()
    colors = {}
    import re
    for line in text.splitlines():
        matches = re.findall(r"#([0-9a-fA-F]{6})", line)
        lower = line.lower()
        for m in matches:
            if "primary" in lower and "primary" not in colors:
                colors["primary"] = f"#{m}"
            elif "secondary" in lower and "secondary" not in colors:
                colors["secondary"] = f"#{m}"
            elif "accent" in lower and "accent" not in colors:
                colors["accent"] = f"#{m}"
            elif "dark" in lower and "dark" not in colors:
                colors["dark"] = f"#{m}"
            elif "light" in lower and "light" not in colors:
                colors["light"] = f"#{m}"
    return colors


# ============================================================
# Theme System
# ============================================================

# Slide dimensions (fixed by Google Slides: 10" x 5.625" at 72 DPI)
SLIDE_WIDTH_PT = 720
SLIDE_HEIGHT_PT = 405

THEME_PROFESSIONAL = {
    "name": "professional",
    "colors": {
        "primary": "#1a1a2e",
        "secondary": "#16213e",
        "accent": "#0f3460",
        "highlight": "#e94560",
        "background": "#ffffff",
        "surface": "#f8f9fa",
        "textPrimary": "#202124",
        "textSecondary": "#5f6368",
        "border": "#dadce0",
        "success": "#34a853",
        "warning": "#fbbc04",
        "error": "#ea4335",
    },
    "fonts": {
        "heading": "Montserrat",
        "body": "Open Sans",
        "mono": "Roboto Mono",
    },
    "typography": {
        "title": {"size": 40, "weight": 700, "lineSpacing": 110},
        "heading1": {"size": 32, "weight": 700, "lineSpacing": 115},
        "heading2": {"size": 24, "weight": 600, "lineSpacing": 120},
        "heading3": {"size": 20, "weight": 600, "lineSpacing": 125},
        "body": {"size": 16, "weight": 400, "lineSpacing": 150},
        "bodyLarge": {"size": 18, "weight": 400, "lineSpacing": 145},
        "caption": {"size": 12, "weight": 400, "lineSpacing": 140},
        "label": {"size": 11, "weight": 600, "lineSpacing": 130},
    },
    "spacing": {
        "slideMargin": 40,
        "contentWidth": 640,
        "elementGap": 16,
        "sectionGap": 32,
    },
}

THEME_CGK = {
    "name": "cgk",
    "colors": {
        "primary": "#000000",
        "secondary": "#2a2a2a",
        "accent": "#B4AA8A",
        "highlight": "#3B3B3B",
        "background": "#FCFCFC",
        "surface": "#FFFFFF",
        "textPrimary": "#2A2A2A",
        "textSecondary": "#3B3B3B",
        "border": "#E0E0E0",
        "success": "#34a853",
        "warning": "#fbbc04",
        "error": "#ea4335",
    },
    "fonts": {
        "heading": "Arial",
        "body": "Arial",
        "mono": "Roboto Mono",
    },
    "typography": {
        "title": {"size": 36, "weight": 700, "lineSpacing": 110},
        "heading1": {"size": 28, "weight": 700, "lineSpacing": 115},
        "heading2": {"size": 22, "weight": 700, "lineSpacing": 120},
        "heading3": {"size": 18, "weight": 700, "lineSpacing": 125},
        "body": {"size": 16, "weight": 400, "lineSpacing": 150},
        "bodyLarge": {"size": 18, "weight": 400, "lineSpacing": 140},
        "caption": {"size": 13, "weight": 400, "lineSpacing": 140},
        "label": {"size": 11, "weight": 700, "lineSpacing": 130},
    },
    "spacing": {
        "slideMargin": 40,
        "contentWidth": 640,
        "elementGap": 16,
        "sectionGap": 32,
    },
}

THEME_RAWDOG = {
    "name": "rawdog",
    "colors": {
        "primary": "#374d42",
        "secondary": "#2d3f35",
        "accent": "#6a947c",
        "highlight": "#C63B32",
        "background": "#FFFFFF",
        "surface": "#f0f4f2",
        "textPrimary": "#000000",
        "textSecondary": "#7a7a7a",
        "border": "#dce7e1",
        "success": "#79B26B",
        "warning": "#fbbc04",
        "error": "#d32f2f",
    },
    "fonts": {
        "heading": "Helvetica Neue",
        "body": "Helvetica Neue",
        "mono": "Roboto Mono",
    },
    "typography": {
        "title": {"size": 40, "weight": 700, "lineSpacing": 110},
        "heading1": {"size": 32, "weight": 700, "lineSpacing": 115},
        "heading2": {"size": 24, "weight": 600, "lineSpacing": 120},
        "heading3": {"size": 20, "weight": 600, "lineSpacing": 125},
        "body": {"size": 16, "weight": 400, "lineSpacing": 150},
        "bodyLarge": {"size": 18, "weight": 400, "lineSpacing": 145},
        "caption": {"size": 12, "weight": 400, "lineSpacing": 140},
        "label": {"size": 11, "weight": 600, "lineSpacing": 130},
    },
    "spacing": {
        "slideMargin": 40,
        "contentWidth": 640,
        "elementGap": 16,
        "sectionGap": 32,
    },
}

THEME_VITAHUSTLE = {
    "name": "vitahustle",
    "colors": {
        "primary": "#0A0A0A",
        "secondary": "#1A1A1A",
        "accent": "#C8E812",
        "highlight": "#7B2FBE",
        "background": "#FFFFFF",
        "surface": "#F5F5F5",
        "textPrimary": "#1A1A1A",
        "textSecondary": "#666666",
        "border": "#E0E0E0",
        "success": "#7AAA3C",
        "warning": "#FFD700",
        "error": "#ea4335",
    },
    "fonts": {
        "heading": "Arial",
        "body": "Arial",
        "mono": "Roboto Mono",
    },
    "typography": {
        "title": {"size": 42, "weight": 700, "lineSpacing": 105},
        "heading1": {"size": 34, "weight": 700, "lineSpacing": 110},
        "heading2": {"size": 26, "weight": 700, "lineSpacing": 115},
        "heading3": {"size": 20, "weight": 700, "lineSpacing": 120},
        "body": {"size": 16, "weight": 400, "lineSpacing": 150},
        "bodyLarge": {"size": 18, "weight": 400, "lineSpacing": 140},
        "caption": {"size": 12, "weight": 400, "lineSpacing": 140},
        "label": {"size": 11, "weight": 700, "lineSpacing": 130},
    },
    "spacing": {
        "slideMargin": 40,
        "contentWidth": 640,
        "elementGap": 16,
        "sectionGap": 32,
    },
}

THEME_DARK = {
    "name": "dark",
    "colors": {
        "primary": "#e8eaed",
        "secondary": "#bdc1c6",
        "accent": "#8ab4f8",
        "highlight": "#f28b82",
        "background": "#1e1e1e",
        "surface": "#2d2d2d",
        "textPrimary": "#e8eaed",
        "textSecondary": "#9aa0a6",
        "border": "#3c4043",
        "success": "#81c995",
        "warning": "#fdd663",
        "error": "#f28b82",
    },
    "fonts": {
        "heading": "Montserrat",
        "body": "Open Sans",
        "mono": "Roboto Mono",
    },
    "typography": THEME_PROFESSIONAL["typography"],
    "spacing": THEME_PROFESSIONAL["spacing"],
}

_theme_cache = None

def _load_brand_theme():
    """Load brand theme from workspace/brand/ files, with smart fallbacks.

    Detection order:
    1. BRAND_NAME env var -> known theme map
    2. Brand files (IDENTITY.md, email-design-system.md, colors.md) -> keyword detection
    3. Profile state dir name -> known profile map
    4. themes/ directory -> custom JSON themes
    5. Fallback -> THEME_PROFESSIONAL
    """
    global _theme_cache
    if _theme_cache is not None:
        return _theme_cache

    ws = _workspace_root()

    # 1. Check BRAND_NAME env var
    brand_name = os.environ.get("BRAND_NAME", "").lower()
    brand_map = {
        "cgk": THEME_CGK, "cgk linens": THEME_CGK,
        "rawdog": THEME_RAWDOG, "justrawdogit": THEME_RAWDOG,
        "vitahustle": THEME_VITAHUSTLE,
    }
    if brand_name in brand_map:
        _theme_cache = brand_map[brand_name]
        return _theme_cache

    # 2. Detect from brand files content
    brand_keywords = {
        "cgk": THEME_CGK, "cgk linens": THEME_CGK,
        "rawdog": THEME_RAWDOG, "justrawdogit": THEME_RAWDOG, "beef tallow": THEME_RAWDOG,
        "vitahustle": THEME_VITAHUSTLE, "vita hustle": THEME_VITAHUSTLE,
    }
    for fname in ("IDENTITY.md", "email-design-system.md", "colors.md", "brand-guide.md"):
        brand_file = ws / "workspace" / "brand" / fname
        if not brand_file.exists():
            brand_file = ws / "workspace" / fname
        if brand_file.exists():
            try:
                text = brand_file.read_text(errors="ignore").lower()
                for keyword, theme in brand_keywords.items():
                    if keyword in text:
                        _theme_cache = theme
                        return _theme_cache
            except OSError:
                pass

    # 3. Detect from profile directory name
    ws_name = str(ws).lower()
    if "rawdog" in ws_name:
        _theme_cache = THEME_RAWDOG
        return _theme_cache
    elif "vitahustle" in ws_name:
        _theme_cache = THEME_VITAHUSTLE
        return _theme_cache
    elif ws_name.endswith(".openclaw"):
        _theme_cache = THEME_CGK
        return _theme_cache

    # 4. Try themes/ directory for custom JSON themes
    themes_dir = SKILL_DIR / "themes"
    if themes_dir.is_dir():
        brand_slug = brand_name.replace(" ", "-")
        for candidate in [brand_slug, "default"]:
            if not candidate:
                continue
            theme_file = themes_dir / f"{candidate}.json"
            if theme_file.exists():
                try:
                    custom = json.loads(theme_file.read_text())
                    merged = {**THEME_PROFESSIONAL, **custom}
                    merged["colors"] = {**THEME_PROFESSIONAL["colors"], **custom.get("colors", {})}
                    merged["fonts"] = {**THEME_PROFESSIONAL["fonts"], **custom.get("fonts", {})}
                    merged["typography"] = {**THEME_PROFESSIONAL["typography"], **custom.get("typography", {})}
                    merged["spacing"] = {**THEME_PROFESSIONAL["spacing"], **custom.get("spacing", {})}
                    _theme_cache = merged
                    return _theme_cache
                except (json.JSONDecodeError, OSError):
                    pass

    # 5. Fallback
    _theme_cache = THEME_PROFESSIONAL
    return _theme_cache


# ============================================================
# AUTH commands
# ============================================================

def cmd_auth_status(_args):
    """Check if credentials exist and are valid."""
    cp = _creds_path()
    if not cp.exists():
        print("Status: NOT CONFIGURED")
        print(f"No credentials found at: {cp}")
        print(f"Run: uv run {SCRIPTS_DIR / 'setup_oauth.py'}")
        return

    try:
        creds_data = json.loads(cp.read_text())
        account = creds_data.get("account", "unknown")

        # Try to actually use the credentials
        creds = _get_credentials()
        if creds.valid:
            print(f"Status: AUTHENTICATED")
            print(f"Account: {account}")
            print(f"Scopes: {', '.join(creds_data.get('scopes', []))}")
            print(f"Credentials: {cp}")
        else:
            print(f"Status: TOKEN EXPIRED (will auto-refresh on next use)")
            print(f"Account: {account}")
    except Exception as e:
        print(f"Status: ERROR")
        print(f"Details: {e}")
        print(f"Try re-running: uv run {SCRIPTS_DIR / 'setup_oauth.py'}")


# ============================================================
# SLIDES commands
# ============================================================

def cmd_slides_create(args):
    """Create a new presentation, optionally removing the default blank slide."""
    svc = _slides_service()
    body = {"title": args.title}
    pres = svc.presentations().create(body=body).execute()
    pres_id = pres["presentationId"]

    # Strip placeholder elements from the default slide ("Click to add title/subtitle")
    if args.clean:
        slides = pres.get("slides", [])
        if slides:
            elements = slides[0].get("pageElements", [])
            if elements:
                del_requests = [{"deleteObject": {"objectId": el["objectId"]}}
                                for el in elements]
                svc.presentations().batchUpdate(
                    presentationId=pres_id, body={"requests": del_requests}
                ).execute()

    _auto_share_org(pres_id)

    url = f"https://docs.google.com/presentation/d/{pres_id}/edit"
    print(json.dumps({"presentationId": pres_id, "url": url, "orgShared": True}, indent=2))


def cmd_slides_get(args):
    """Get presentation metadata and structure."""
    svc = _slides_service()
    pres = svc.presentations().get(presentationId=args.presentation_id).execute()
    # Summarize structure
    summary = {
        "presentationId": pres["presentationId"],
        "title": pres.get("title", ""),
        "url": f"https://docs.google.com/presentation/d/{pres['presentationId']}/edit",
        "slideCount": len(pres.get("slides", [])),
        "slides": [],
    }
    for i, slide in enumerate(pres.get("slides", [])):
        slide_info = {
            "index": i,
            "objectId": slide["objectId"],
            "elementCount": len(slide.get("pageElements", [])),
        }
        summary["slides"].append(slide_info)
    print(json.dumps(summary, indent=2))


def _resolve_slide_id(svc, pres_id, slide_idx):
    """Resolve a slide index to its objectId."""
    pres = svc.presentations().get(presentationId=pres_id).execute()
    slides = pres.get("slides", [])
    if slide_idx < 0 or slide_idx >= len(slides):
        print(f"ERROR: Slide index {slide_idx} out of range (0-{len(slides)-1}).", file=sys.stderr)
        sys.exit(1)
    return slides[slide_idx]["objectId"]


def cmd_slides_add_slide(args):
    """Add a new slide with optional layout."""
    svc = _slides_service()
    slide_id = _gen_id()
    request = {
        "createSlide": {
            "objectId": slide_id,
            "slideLayoutReference": {"predefinedLayout": args.layout},
        }
    }
    if args.index is not None:
        request["createSlide"]["insertionIndex"] = args.index

    svc.presentations().batchUpdate(
        presentationId=args.presentation_id, body={"requests": [request]}
    ).execute()
    print(json.dumps({"slideId": slide_id, "layout": args.layout}))


def cmd_slides_add_text(args):
    """Add a styled text box to a slide with theme-aware defaults."""
    svc = _slides_service()
    slide_id = _resolve_slide_id(svc, args.presentation_id, args.slide_index)
    style = _parse_json_arg(args.style, "style") if args.style else {}
    theme = _load_brand_theme()
    sp = theme["spacing"]

    # Resolve heading level
    heading = style.get("heading") or args.heading
    if heading and heading in theme["typography"]:
        typo = theme["typography"][heading]
    elif heading == "title":
        typo = theme["typography"]["title"]
    elif heading == "subtitle":
        typo = theme["typography"].get("bodyLarge", theme["typography"]["body"])
    else:
        typo = theme["typography"]["body"]

    is_heading = heading in ("heading1", "heading2", "heading3", "title")

    # Apply theme defaults, allow style overrides
    x = style.get("x", sp["slideMargin"])
    y = style.get("y", sp["slideMargin"])
    width = style.get("width", sp["contentWidth"])
    height = style.get("height", 60)
    font_size = style.get("fontSize", typo["size"])
    bold = style.get("bold", typo["weight"] >= 600)
    italic = style.get("italic", False)
    color = style.get("color", theme["colors"]["textPrimary"])
    alignment = style.get("alignment", "START")
    font_family = style.get("fontFamily",
                            theme["fonts"]["heading"] if is_heading else theme["fonts"]["body"])
    line_spacing = style.get("lineSpacing", typo.get("lineSpacing", 150))
    space_above = style.get("spaceAbove", 0)
    space_below = style.get("spaceBelow", 0)

    textbox_id = _gen_id()
    text_style = {
        "bold": bold,
        "italic": italic,
        "fontSize": _pt(font_size),
        "foregroundColor": _rgb_color(color),
        "fontFamily": font_family,
    }
    text_fields = "bold,italic,fontSize,foregroundColor,fontFamily"

    # Add weighted font family for precise weight control
    if typo["weight"] != 400:
        text_style["weightedFontFamily"] = {
            "fontFamily": font_family,
            "weight": typo["weight"],
        }
        text_fields += ",weightedFontFamily"

    requests = [
        {
            "createShape": {
                "objectId": textbox_id,
                "shapeType": "TEXT_BOX",
                "elementProperties": {
                    "pageObjectId": slide_id,
                    "size": {"height": _pt(height), "width": _pt(width)},
                    "transform": {
                        "scaleX": 1, "scaleY": 1,
                        "translateX": x, "translateY": y, "unit": "PT",
                    },
                },
            }
        },
        {"insertText": {"objectId": textbox_id, "insertionIndex": 0, "text": args.text}},
        {
            "updateTextStyle": {
                "objectId": textbox_id,
                "textRange": {"type": "ALL"},
                "style": text_style,
                "fields": text_fields,
            }
        },
        {
            "updateParagraphStyle": {
                "objectId": textbox_id,
                "textRange": {"type": "ALL"},
                "style": {
                    "alignment": alignment,
                    "lineSpacing": line_spacing,
                    "spaceAbove": _pt(space_above),
                    "spaceBelow": _pt(space_below),
                },
                "fields": "alignment,lineSpacing,spaceAbove,spaceBelow",
            }
        },
    ]

    svc.presentations().batchUpdate(
        presentationId=args.presentation_id, body={"requests": requests}
    ).execute()
    print(json.dumps({"textBoxId": textbox_id, "slideIndex": args.slide_index}))


def cmd_slides_add_table(args):
    """Add a professionally styled table to a slide with theme-aware defaults."""
    svc = _slides_service()
    slide_id = _resolve_slide_id(svc, args.presentation_id, args.slide_index)
    data = _parse_json_arg(args.data, "data")
    style = _parse_json_arg(args.style, "style") if args.style else {}
    theme = _load_brand_theme()

    if not data or not isinstance(data, list) or not isinstance(data[0], list):
        print("ERROR: --data must be a 2D JSON array: [[\"H1\",\"H2\"],[\"v1\",\"v2\"]]", file=sys.stderr)
        sys.exit(1)

    rows = len(data)
    cols = len(data[0])
    table_id = _gen_id()

    header_color = style.get("headerColor", theme["colors"]["primary"])
    header_text_color = style.get("headerTextColor", "#ffffff")
    zebra = style.get("zebra", True)
    zebra_color = style.get("zebraColor", theme["colors"]["surface"])
    border_color = style.get("borderColor", theme["colors"]["border"])
    font_size = style.get("fontSize", theme["typography"]["caption"]["size"])
    header_font_size = style.get("headerFontSize", font_size)
    font_family = style.get("fontFamily", theme["fonts"]["body"])
    content_align = style.get("contentAlignment", "MIDDLE")

    requests = [
        {
            "createTable": {
                "objectId": table_id,
                "elementProperties": {"pageObjectId": slide_id},
                "rows": rows,
                "columns": cols,
            }
        }
    ]

    # Insert text into each cell
    for r_idx, row in enumerate(data):
        for c_idx, cell_val in enumerate(row):
            requests.append({
                "insertText": {
                    "objectId": table_id,
                    "cellLocation": {"rowIndex": r_idx, "columnIndex": c_idx},
                    "insertionIndex": 0,
                    "text": str(cell_val),
                }
            })

    # Style header row background
    requests.append({
        "updateTableCellProperties": {
            "objectId": table_id,
            "tableRange": {
                "location": {"rowIndex": 0, "columnIndex": 0},
                "rowSpan": 1, "columnSpan": cols,
            },
            "tableCellProperties": {
                "tableCellBackgroundFill": _solid_fill(header_color),
                "contentAlignment": "MIDDLE",
            },
            "fields": "tableCellBackgroundFill,contentAlignment",
        }
    })

    # Style header text (bold, white, centered, font family)
    for c_idx in range(cols):
        requests.append({
            "updateTextStyle": {
                "objectId": table_id,
                "cellLocation": {"rowIndex": 0, "columnIndex": c_idx},
                "textRange": {"type": "ALL"},
                "style": {
                    "bold": True,
                    "fontSize": _pt(header_font_size),
                    "foregroundColor": _rgb_color(header_text_color),
                    "fontFamily": font_family,
                },
                "fields": "bold,fontSize,foregroundColor,fontFamily",
            }
        })
        requests.append({
            "updateParagraphStyle": {
                "objectId": table_id,
                "cellLocation": {"rowIndex": 0, "columnIndex": c_idx},
                "textRange": {"type": "ALL"},
                "style": {"alignment": "CENTER"},
                "fields": "alignment",
            }
        })

    # Zebra striping + content alignment on data rows
    for r_idx in range(1, rows):
        cell_props = {"contentAlignment": content_align}
        fields = "contentAlignment"
        if zebra and r_idx % 2 == 0:
            cell_props["tableCellBackgroundFill"] = _solid_fill(zebra_color)
            fields += ",tableCellBackgroundFill"
        requests.append({
            "updateTableCellProperties": {
                "objectId": table_id,
                "tableRange": {
                    "location": {"rowIndex": r_idx, "columnIndex": 0},
                    "rowSpan": 1, "columnSpan": cols,
                },
                "tableCellProperties": cell_props,
                "fields": fields,
            }
        })

    # Style all data cell text with font family
    for r_idx in range(1, rows):
        for c_idx in range(cols):
            requests.append({
                "updateTextStyle": {
                    "objectId": table_id,
                    "cellLocation": {"rowIndex": r_idx, "columnIndex": c_idx},
                    "textRange": {"type": "ALL"},
                    "style": {
                        "fontSize": _pt(font_size),
                        "foregroundColor": _rgb_color(theme["colors"]["textPrimary"]),
                        "fontFamily": font_family,
                    },
                    "fields": "fontSize,foregroundColor,fontFamily",
                }
            })

    # Table borders
    requests.append({
        "updateTableBorderProperties": {
            "objectId": table_id,
            "borderPosition": "ALL",
            "tableBorderProperties": {
                "tableBorderFill": _solid_fill(border_color),
                "weight": _pt(0.5),
                "dashStyle": "SOLID",
            },
            "fields": "tableBorderFill,weight,dashStyle",
        }
    })

    # Column widths if specified
    col_widths = style.get("columnWidths")
    if col_widths:
        for i, w in enumerate(col_widths):
            if i < cols:
                requests.append({
                    "updateTableColumnProperties": {
                        "objectId": table_id,
                        "columnIndices": [i],
                        "tableColumnProperties": {"columnWidth": _pt(w)},
                        "fields": "columnWidth",
                    }
                })

    # Header row minimum height
    header_row_height = style.get("headerRowHeight", 36)
    requests.append({
        "updateTableRowProperties": {
            "objectId": table_id,
            "rowIndices": [0],
            "tableRowProperties": {"minRowHeight": _pt(header_row_height)},
            "fields": "minRowHeight",
        }
    })

    svc.presentations().batchUpdate(
        presentationId=args.presentation_id, body={"requests": requests}
    ).execute()
    print(json.dumps({"tableId": table_id, "rows": rows, "cols": cols}))


def cmd_slides_add_image(args):
    """Add an image to a slide from a public URL."""
    svc = _slides_service()
    slide_id = _resolve_slide_id(svc, args.presentation_id, args.slide_index)
    image_id = _gen_id()

    x = args.x if args.x is not None else _emu_from_inches(1)
    y = args.y if args.y is not None else _emu_from_inches(1)
    width = args.width if args.width is not None else _emu_from_inches(4)
    height = args.height if args.height is not None else _emu_from_inches(3)

    request = {
        "createImage": {
            "objectId": image_id,
            "url": args.url,
            "elementProperties": {
                "pageObjectId": slide_id,
                "size": {
                    "height": _emu(height),
                    "width": _emu(width),
                },
                "transform": {
                    "scaleX": 1, "scaleY": 1,
                    "translateX": x, "translateY": y,
                    "unit": "EMU",
                },
            },
        }
    }

    svc.presentations().batchUpdate(
        presentationId=args.presentation_id, body={"requests": [request]}
    ).execute()
    print(json.dumps({"imageId": image_id}))


def cmd_slides_add_shape(args):
    """Add a shape to a slide."""
    svc = _slides_service()
    slide_id = _resolve_slide_id(svc, args.presentation_id, args.slide_index)
    style = _parse_json_arg(args.style, "style") if args.style else {}

    shape_id = _gen_id()
    x = style.get("x", 100)
    y = style.get("y", 100)
    width = style.get("width", 200)
    height = style.get("height", 100)
    fill_color = style.get("fillColor")
    outline_color = style.get("outlineColor")

    requests = [
        {
            "createShape": {
                "objectId": shape_id,
                "shapeType": args.shape_type,
                "elementProperties": {
                    "pageObjectId": slide_id,
                    "size": {"height": _pt(height), "width": _pt(width)},
                    "transform": {
                        "scaleX": 1, "scaleY": 1,
                        "translateX": x, "translateY": y, "unit": "PT",
                    },
                },
            }
        }
    ]

    # Apply fill/outline if specified
    props = {}
    fields = []
    if fill_color:
        props["shapeBackgroundFill"] = _solid_fill(fill_color)
        fields.append("shapeBackgroundFill")
    if outline_color:
        props["outline"] = {
            "outlineFill": _solid_fill(outline_color),
            "weight": _pt(style.get("outlineWeight", 1)),
            "dashStyle": style.get("dashStyle", "SOLID"),
        }
        fields.append("outline")
    if fields:
        requests.append({
            "updateShapeProperties": {
                "objectId": shape_id,
                "shapeProperties": props,
                "fields": ",".join(fields),
            }
        })

    svc.presentations().batchUpdate(
        presentationId=args.presentation_id, body={"requests": requests}
    ).execute()
    print(json.dumps({"shapeId": shape_id, "type": args.shape_type}))


def cmd_slides_set_background(args):
    """Set a slide's background to a solid color or image."""
    svc = _slides_service()
    slide_id = _resolve_slide_id(svc, args.presentation_id, args.slide_index)

    page_props = {}
    if args.color:
        page_props["pageBackgroundFill"] = _solid_fill(args.color)
    elif args.image:
        page_props["pageBackgroundFill"] = {
            "stretchedPictureFill": {"contentUrl": args.image}
        }
    else:
        print("ERROR: Provide --color or --image.", file=sys.stderr)
        sys.exit(1)

    svc.presentations().batchUpdate(
        presentationId=args.presentation_id,
        body={"requests": [{
            "updatePageProperties": {
                "objectId": slide_id,
                "pageProperties": page_props,
                "fields": "pageBackgroundFill",
            }
        }]},
    ).execute()
    print(json.dumps({"slideIndex": args.slide_index, "background": "set"}))


def cmd_slides_apply_theme(args):
    """Apply brand theme to presentation: master color scheme + slide backgrounds."""
    svc = _slides_service()
    theme = _load_brand_theme()
    colors = theme["colors"]

    pres = svc.presentations().get(presentationId=args.presentation_id).execute()
    slides = pres.get("slides", [])
    masters = pres.get("masters", [])

    if not slides:
        print("No slides to theme.")
        return

    requests = []

    # Update master color scheme (all 12 ThemeColorTypes)
    if masters:
        master_id = masters[0]["objectId"]
        color_scheme = [
            {"type": "DARK1", "color": _hex_to_rgb(colors["textPrimary"])},
            {"type": "LIGHT1", "color": _hex_to_rgb(colors["background"])},
            {"type": "DARK2", "color": _hex_to_rgb(colors["secondary"])},
            {"type": "LIGHT2", "color": _hex_to_rgb(colors["surface"])},
            {"type": "ACCENT1", "color": _hex_to_rgb(colors["primary"])},
            {"type": "ACCENT2", "color": _hex_to_rgb(colors["accent"])},
            {"type": "ACCENT3", "color": _hex_to_rgb(colors["success"])},
            {"type": "ACCENT4", "color": _hex_to_rgb(colors["warning"])},
            {"type": "ACCENT5", "color": _hex_to_rgb(colors["error"])},
            {"type": "ACCENT6", "color": _hex_to_rgb(colors["highlight"])},
            {"type": "HYPERLINK", "color": _hex_to_rgb(colors["accent"])},
            {"type": "FOLLOWED_HYPERLINK", "color": _hex_to_rgb(colors["textSecondary"])},
        ]
        requests.append({
            "updatePageProperties": {
                "objectId": master_id,
                "pageProperties": {"colorScheme": {"colors": color_scheme}},
                "fields": "colorScheme",
            }
        })

    # Only set backgrounds on slides that don't already have explicit backgrounds
    # (i.e., skip slides created by add-layout which handle their own backgrounds)
    if args.set_backgrounds:
        for i, slide in enumerate(slides):
            bg_color = colors["primary"] if i == 0 else colors["background"]
            requests.append({
                "updatePageProperties": {
                    "objectId": slide["objectId"],
                    "pageProperties": {"pageBackgroundFill": _solid_fill(bg_color)},
                    "fields": "pageBackgroundFill",
                }
            })

    svc.presentations().batchUpdate(
        presentationId=args.presentation_id, body={"requests": requests}
    ).execute()
    print(json.dumps({
        "themed": True,
        "themeName": theme["name"],
        "masterUpdated": bool(masters),
        "slidesAffected": len(slides),
    }))


def cmd_slides_add_chart(args):
    """Add a chart to a slide (matplotlib image or Sheets-linked)."""
    svc = _slides_service()
    slide_id = _resolve_slide_id(svc, args.presentation_id, args.slide_index)
    data = _parse_json_arg(args.data, "data")

    if args.chart_source == "sheets":
        _slides_add_chart_sheets(svc, args, slide_id, data)
    else:
        _slides_add_chart_image(svc, args, slide_id, data)


def _slides_add_chart_image(svc, args, slide_id, data):
    """Generate a professional chart image with matplotlib and insert into slide."""
    import matplotlib
    matplotlib.use("Agg")
    import matplotlib.pyplot as plt

    labels = data.get("labels", [])
    datasets = data.get("datasets", [])
    if not datasets:
        print("ERROR: --data must include 'datasets' array.", file=sys.stderr)
        sys.exit(1)

    theme = _load_brand_theme()
    tc = theme["colors"]
    chart_colors = [
        tc["primary"], tc["accent"], tc["success"],
        tc["warning"], tc["error"], "#46bdc6", "#7b1fa2",
    ]

    # Professional matplotlib styling
    plt.rcParams.update({
        "font.family": "sans-serif",
        "font.sans-serif": [theme["fonts"]["body"], "Arial", "Helvetica"],
        "font.size": 11,
        "axes.titlesize": 14,
        "axes.titleweight": "bold",
        "axes.labelsize": 11,
        "axes.facecolor": "#ffffff",
        "axes.edgecolor": tc["border"],
        "axes.grid": True,
        "axes.axisbelow": True,
        "grid.color": "#f0f0f0",
        "grid.linewidth": 0.5,
        "xtick.color": tc["textSecondary"],
        "ytick.color": tc["textSecondary"],
        "text.color": tc["textPrimary"],
        "figure.facecolor": "#ffffff",
        "legend.frameon": False,
        "legend.fontsize": 10,
    })

    fig, ax = plt.subplots(figsize=(9, 5.5))
    chart_type = args.chart_type

    for i, ds in enumerate(datasets):
        color = chart_colors[i % len(chart_colors)]
        values = ds.get("values", [])
        label = ds.get("label", f"Series {i+1}")

        if chart_type == "bar":
            import numpy as np
            x = np.arange(len(labels))
            w = 0.8 / len(datasets)
            ax.bar(x + i * w - 0.4 + w / 2, values, w, label=label, color=color,
                   edgecolor="white", linewidth=0.5)
            ax.set_xticks(x)
            ax.set_xticklabels(labels, rotation=45, ha="right")
        elif chart_type == "line":
            ax.plot(labels, values, marker="o", label=label, color=color,
                    linewidth=2.5, markersize=6)
        elif chart_type == "pie":
            ax.pie(values, labels=labels, autopct="%1.1f%%",
                   colors=chart_colors[:len(values)], startangle=90,
                   wedgeprops={"edgecolor": "white", "linewidth": 2})
            ax.set_aspect("equal")
        elif chart_type in ("column", "area"):
            if chart_type == "area":
                ax.fill_between(range(len(values)), values, alpha=0.15, color=color)
            ax.plot(range(len(values)), values, marker="o", label=label, color=color,
                    linewidth=2.5, markersize=6)
            ax.set_xticks(range(len(labels)))
            ax.set_xticklabels(labels, rotation=45, ha="right")
        elif chart_type == "scatter":
            x_vals = ds.get("x", list(range(len(values))))
            ax.scatter(x_vals, values, label=label, color=color, s=60, edgecolors="white")

    if chart_type != "pie":
        ax.spines["top"].set_visible(False)
        ax.spines["right"].set_visible(False)
        ax.spines["left"].set_linewidth(0.5)
        ax.spines["bottom"].set_linewidth(0.5)
        ax.legend()

    title = data.get("title", "")
    if title:
        ax.set_title(title, fontsize=14, fontweight="bold", pad=12)
    x_title = data.get("xAxisTitle", "")
    y_title = data.get("yAxisTitle", "")
    if x_title:
        ax.set_xlabel(x_title)
    if y_title:
        ax.set_ylabel(y_title)
    plt.tight_layout()

    # Save to temp file, upload to Drive, insert into Slides
    with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as tmp:
        fig.savefig(tmp.name, dpi=150, bbox_inches="tight", facecolor="white")
        plt.close(fig)
        tmp_path = tmp.name

    try:
        drive = _drive_service()
        from googleapiclient.http import MediaFileUpload
        media = MediaFileUpload(tmp_path, mimetype="image/png")
        uploaded = drive.files().create(
            body={"name": f"chart-{_gen_id()}.png", "mimeType": "image/png"},
            media_body=media,
            fields="id,webContentLink",
        ).execute()
        file_id = uploaded["id"]

        # Make publicly readable for Slides to fetch
        drive.permissions().create(
            fileId=file_id,
            body={"type": "anyone", "role": "reader"},
        ).execute()

        link = f"https://drive.google.com/uc?id={file_id}"

        image_id = _gen_id()
        svc.presentations().batchUpdate(
            presentationId=args.presentation_id,
            body={"requests": [{
                "createImage": {
                    "objectId": image_id,
                    "url": link,
                    "elementProperties": {
                        "pageObjectId": slide_id,
                        "size": {
                            "width": _emu(_emu_from_inches(7)),
                            "height": _emu(_emu_from_inches(4.5)),
                        },
                        "transform": {
                            "scaleX": 1, "scaleY": 1,
                            "translateX": _emu_from_inches(1.5),
                            "translateY": _emu_from_inches(0.5),
                            "unit": "EMU",
                        },
                    },
                }
            }]},
        ).execute()
        print(json.dumps({"chartImageId": image_id, "driveFileId": file_id}))
    finally:
        os.unlink(tmp_path)


def _slides_add_chart_sheets(svc, args, slide_id, data):
    """Create a chart in Sheets and embed it into Slides."""
    sheets_svc = _sheets_service()
    drive = _drive_service()

    labels = data.get("labels", [])
    datasets = data.get("datasets", [])

    # Build spreadsheet data: header row + data rows
    header = ["Category"] + [ds.get("label", f"Series {i+1}") for i, ds in enumerate(datasets)]
    rows = []
    for i, label in enumerate(labels):
        row = [label] + [ds["values"][i] if i < len(ds["values"]) else "" for ds in datasets]
        rows.append(row)

    sheet_data = [header] + rows

    # Create a temporary spreadsheet
    ss = sheets_svc.spreadsheets().create(body={
        "properties": {"title": f"Chart Data - {_gen_id()[:8]}"},
        "sheets": [{"properties": {"title": "Data"}}],
    }).execute()
    ss_id = ss["spreadsheetId"]
    sheet_id = ss["sheets"][0]["properties"]["sheetId"]

    # Write data
    sheets_svc.spreadsheets().values().update(
        spreadsheetId=ss_id,
        range="Data!A1",
        valueInputOption="RAW",
        body={"values": sheet_data},
    ).execute()

    # Map chart type
    chart_type_map = {
        "bar": "BAR", "line": "LINE", "pie": "PIE",
        "column": "COLUMN", "area": "AREA", "scatter": "SCATTER",
    }
    sheets_chart_type = chart_type_map.get(args.chart_type, "BAR")

    # Create chart in the spreadsheet
    end_col = len(header) - 1
    end_row = len(sheet_data)

    chart_spec = {
        "title": data.get("title", ""),
        "basicChart": {
            "chartType": sheets_chart_type,
            "legendPosition": "BOTTOM_LEGEND",
            "axis": [
                {"position": "BOTTOM_AXIS", "title": data.get("xAxisTitle", "")},
                {"position": "LEFT_AXIS", "title": data.get("yAxisTitle", "")},
            ],
            "domains": [{
                "domain": {
                    "sourceRange": {
                        "sources": [{
                            "sheetId": sheet_id,
                            "startRowIndex": 0, "endRowIndex": end_row,
                            "startColumnIndex": 0, "endColumnIndex": 1,
                        }]
                    }
                }
            }],
            "series": [
                {
                    "series": {
                        "sourceRange": {
                            "sources": [{
                                "sheetId": sheet_id,
                                "startRowIndex": 0, "endRowIndex": end_row,
                                "startColumnIndex": col_idx + 1, "endColumnIndex": col_idx + 2,
                            }]
                        }
                    },
                    "targetAxis": "BOTTOM_AXIS" if sheets_chart_type == "BAR" else "LEFT_AXIS",
                }
                for col_idx in range(len(datasets))
            ],
            "headerCount": 1,
        },
    }

    # Use pie chart spec if pie
    if sheets_chart_type == "PIE":
        chart_spec = {
            "title": data.get("title", ""),
            "pieChart": {
                "legendPosition": "RIGHT_LEGEND",
                "domain": {
                    "sourceRange": {
                        "sources": [{
                            "sheetId": sheet_id,
                            "startRowIndex": 0, "endRowIndex": end_row,
                            "startColumnIndex": 0, "endColumnIndex": 1,
                        }]
                    }
                },
                "series": {
                    "sourceRange": {
                        "sources": [{
                            "sheetId": sheet_id,
                            "startRowIndex": 0, "endRowIndex": end_row,
                            "startColumnIndex": 1, "endColumnIndex": 2,
                        }]
                    }
                },
            },
        }

    add_chart_req = {
        "addChart": {
            "chart": {
                "spec": chart_spec,
                "position": {
                    "overlayPosition": {
                        "anchorCell": {"sheetId": sheet_id, "rowIndex": end_row + 1, "columnIndex": 0},
                        "widthPixels": 600,
                        "heightPixels": 400,
                    }
                },
            }
        }
    }

    resp = sheets_svc.spreadsheets().batchUpdate(
        spreadsheetId=ss_id, body={"requests": [add_chart_req]}
    ).execute()

    chart_id = resp["replies"][0]["addChart"]["chart"]["chartId"]

    # Embed chart into Slides
    chart_obj_id = _gen_id()
    svc.presentations().batchUpdate(
        presentationId=args.presentation_id,
        body={"requests": [{
            "createSheetsChart": {
                "objectId": chart_obj_id,
                "spreadsheetId": ss_id,
                "chartId": chart_id,
                "linkingMode": "NOT_LINKED_IMAGE",
                "elementProperties": {
                    "pageObjectId": slide_id,
                    "size": {
                        "width": _emu(_emu_from_inches(7)),
                        "height": _emu(_emu_from_inches(4.5)),
                    },
                    "transform": {
                        "scaleX": 1, "scaleY": 1,
                        "translateX": _emu_from_inches(1.5),
                        "translateY": _emu_from_inches(0.5),
                        "unit": "EMU",
                    },
                },
            }
        }]},
    ).execute()
    print(json.dumps({
        "chartObjectId": chart_obj_id,
        "spreadsheetId": ss_id,
        "sheetsChartId": chart_id,
    }))


def _build_accent_bar(slide_id, theme, position="top"):
    """Build requests for a thin colored accent bar."""
    colors = theme["colors"]
    bar_id = _gen_id()
    bar_h = 4
    if position == "top":
        x, y, w, h = 0, 0, SLIDE_WIDTH_PT, bar_h
    elif position == "bottom":
        x, y, w, h = 0, SLIDE_HEIGHT_PT - bar_h, SLIDE_WIDTH_PT, bar_h
    else:
        x, y, w, h = 0, 0, bar_h, SLIDE_HEIGHT_PT
    return [
        {
            "createShape": {
                "objectId": bar_id,
                "shapeType": "RECTANGLE",
                "elementProperties": {
                    "pageObjectId": slide_id,
                    "size": {"height": _pt(h), "width": _pt(w)},
                    "transform": {"scaleX": 1, "scaleY": 1,
                                  "translateX": x, "translateY": y, "unit": "PT"},
                },
            }
        },
        {
            "updateShapeProperties": {
                "objectId": bar_id,
                "shapeProperties": {
                    "shapeBackgroundFill": _solid_fill(colors["accent"]),
                    "outline": {"propertyState": "NOT_RENDERED"},
                },
                "fields": "shapeBackgroundFill,outline",
            }
        },
    ]


def _build_kpi_card(slide_id, theme, value, label, x, y, width=145, height=95):
    """Build requests for a KPI metric card."""
    colors = theme["colors"]
    fonts = theme["fonts"]
    card_id = _gen_id()
    val_id = _gen_id()
    lbl_id = _gen_id()

    return [
        # Card background
        {
            "createShape": {
                "objectId": card_id,
                "shapeType": "ROUND_RECTANGLE",
                "elementProperties": {
                    "pageObjectId": slide_id,
                    "size": {"height": _pt(height), "width": _pt(width)},
                    "transform": {"scaleX": 1, "scaleY": 1,
                                  "translateX": x, "translateY": y, "unit": "PT"},
                },
            }
        },
        {
            "updateShapeProperties": {
                "objectId": card_id,
                "shapeProperties": {
                    "shapeBackgroundFill": _solid_fill(colors["surface"]),
                    "outline": {
                        "outlineFill": _solid_fill(colors["border"]),
                        "weight": _pt(0.5), "dashStyle": "SOLID",
                    },
                },
                "fields": "shapeBackgroundFill,outline",
            }
        },
        # Value
        {
            "createShape": {
                "objectId": val_id,
                "shapeType": "TEXT_BOX",
                "elementProperties": {
                    "pageObjectId": slide_id,
                    "size": {"height": _pt(height * 0.55), "width": _pt(width)},
                    "transform": {"scaleX": 1, "scaleY": 1,
                                  "translateX": x, "translateY": y + 8, "unit": "PT"},
                },
            }
        },
        {"insertText": {"objectId": val_id, "insertionIndex": 0, "text": str(value)}},
        {
            "updateTextStyle": {
                "objectId": val_id, "textRange": {"type": "ALL"},
                "style": {
                    "bold": True, "fontSize": _pt(28),
                    "foregroundColor": _rgb_color(colors["primary"]),
                    "fontFamily": fonts["heading"],
                },
                "fields": "bold,fontSize,foregroundColor,fontFamily",
            }
        },
        {
            "updateParagraphStyle": {
                "objectId": val_id, "textRange": {"type": "ALL"},
                "style": {"alignment": "CENTER"},
                "fields": "alignment",
            }
        },
        # Label
        {
            "createShape": {
                "objectId": lbl_id,
                "shapeType": "TEXT_BOX",
                "elementProperties": {
                    "pageObjectId": slide_id,
                    "size": {"height": _pt(height * 0.3), "width": _pt(width)},
                    "transform": {"scaleX": 1, "scaleY": 1,
                                  "translateX": x, "translateY": y + height * 0.6, "unit": "PT"},
                },
            }
        },
        {"insertText": {"objectId": lbl_id, "insertionIndex": 0, "text": label}},
        {
            "updateTextStyle": {
                "objectId": lbl_id, "textRange": {"type": "ALL"},
                "style": {
                    "fontSize": _pt(11),
                    "foregroundColor": _rgb_color(colors["textSecondary"]),
                    "fontFamily": fonts["body"],
                },
                "fields": "fontSize,foregroundColor,fontFamily",
            }
        },
        {
            "updateParagraphStyle": {
                "objectId": lbl_id, "textRange": {"type": "ALL"},
                "style": {"alignment": "CENTER"},
                "fields": "alignment",
            }
        },
    ]


def _build_text_element(slide_id, theme, text, style_name, x, y, width, height,
                        color_override=None, align="START"):
    """Build requests for a themed text box."""
    fonts = theme["fonts"]
    typo = theme["typography"].get(style_name, theme["typography"]["body"])
    is_heading = style_name in ("title", "heading1", "heading2", "heading3")
    font = fonts["heading"] if is_heading else fonts["body"]
    color = color_override or theme["colors"]["textPrimary"]
    tid = _gen_id()

    return [
        {
            "createShape": {
                "objectId": tid,
                "shapeType": "TEXT_BOX",
                "elementProperties": {
                    "pageObjectId": slide_id,
                    "size": {"height": _pt(height), "width": _pt(width)},
                    "transform": {"scaleX": 1, "scaleY": 1,
                                  "translateX": x, "translateY": y, "unit": "PT"},
                },
            }
        },
        {"insertText": {"objectId": tid, "insertionIndex": 0, "text": text}},
        {
            "updateTextStyle": {
                "objectId": tid, "textRange": {"type": "ALL"},
                "style": {
                    "bold": typo["weight"] >= 600,
                    "fontSize": _pt(typo["size"]),
                    "foregroundColor": _rgb_color(color),
                    "fontFamily": font,
                },
                "fields": "bold,fontSize,foregroundColor,fontFamily",
            }
        },
        {
            "updateParagraphStyle": {
                "objectId": tid, "textRange": {"type": "ALL"},
                "style": {"alignment": align, "lineSpacing": typo.get("lineSpacing", 150)},
                "fields": "alignment,lineSpacing",
            }
        },
    ]


def cmd_slides_add_layout(args):
    """Add a pre-designed professional slide layout.

    Layout types:
      title          — Dark title slide with accent bar
      section        — Section divider slide
      kpi            — KPI dashboard with 4 metric cards (requires --data with kpis array)
      two-column     — Split layout with left title and right content
      closing        — Thank you / closing slide
      content        — Standard content slide with heading
      ad-card        — Ad creative card with image + 8 metrics + preview link (1 ad per slide)
    """
    svc = _slides_service()
    theme = _load_brand_theme()
    colors = theme["colors"]
    sp = theme["spacing"]
    data = _parse_json_arg(args.data, "data") if args.data else {}

    slide_id = _gen_id()
    requests = [{
        "createSlide": {
            "objectId": slide_id,
            "slideLayoutReference": {"predefinedLayout": "BLANK"},
        }
    }]

    layout = args.layout_type

    if layout == "title":
        title = data.get("title", "Presentation Title")
        subtitle = data.get("subtitle", "")
        # Dark background
        requests.append({
            "updatePageProperties": {
                "objectId": slide_id,
                "pageProperties": {"pageBackgroundFill": _solid_fill(colors["primary"])},
                "fields": "pageBackgroundFill",
            }
        })
        # Accent bar at top
        requests.extend(_build_accent_bar(slide_id, theme, "top"))
        # Title
        requests.extend(_build_text_element(
            slide_id, theme, title, "title",
            sp["slideMargin"], SLIDE_HEIGHT_PT * 0.3, sp["contentWidth"], 60,
            color_override="#ffffff", align="CENTER"))
        # Subtitle
        if subtitle:
            requests.extend(_build_text_element(
                slide_id, theme, subtitle, "bodyLarge",
                sp["slideMargin"], SLIDE_HEIGHT_PT * 0.3 + 65, sp["contentWidth"], 40,
                color_override=colors["textSecondary"] if colors["primary"] == "#000000" else "#cccccc",
                align="CENTER"))

    elif layout == "section":
        title = data.get("title", "Section Title")
        # Dark background
        requests.append({
            "updatePageProperties": {
                "objectId": slide_id,
                "pageProperties": {"pageBackgroundFill": _solid_fill(colors["secondary"])},
                "fields": "pageBackgroundFill",
            }
        })
        requests.extend(_build_accent_bar(slide_id, theme, "bottom"))
        requests.extend(_build_text_element(
            slide_id, theme, title, "title",
            sp["slideMargin"], SLIDE_HEIGHT_PT * 0.35, sp["contentWidth"], 60,
            color_override="#ffffff", align="CENTER"))

    elif layout == "kpi":
        title = data.get("title", "Key Metrics")
        kpis = data.get("kpis", [])
        # Light background
        requests.append({
            "updatePageProperties": {
                "objectId": slide_id,
                "pageProperties": {"pageBackgroundFill": _solid_fill(colors["background"])},
                "fields": "pageBackgroundFill",
            }
        })
        requests.extend(_build_accent_bar(slide_id, theme, "top"))
        # Title
        requests.extend(_build_text_element(
            slide_id, theme, title, "heading1",
            sp["slideMargin"], 20, sp["contentWidth"], 45, align="START"))
        # KPI cards in a row
        n = min(len(kpis), 4) if kpis else 4
        if n == 0:
            n = 4
            kpis = [{"value": "—", "label": f"Metric {i+1}"} for i in range(4)]
        card_gap = 16
        total_gap = card_gap * (n - 1)
        card_w = (sp["contentWidth"] - total_gap) / n
        card_h = 100
        card_y = 85
        for i in range(n):
            kpi = kpis[i] if i < len(kpis) else {"value": "—", "label": "—"}
            card_x = sp["slideMargin"] + i * (card_w + card_gap)
            requests.extend(_build_kpi_card(
                slide_id, theme, kpi.get("value", "—"), kpi.get("label", "—"),
                card_x, card_y, card_w, card_h))

    elif layout == "two-column":
        left_title = data.get("title", "")
        left_body = data.get("body", "")
        right_body = data.get("rightBody", "")
        # Light background
        requests.append({
            "updatePageProperties": {
                "objectId": slide_id,
                "pageProperties": {"pageBackgroundFill": _solid_fill(colors["background"])},
                "fields": "pageBackgroundFill",
            }
        })
        requests.extend(_build_accent_bar(slide_id, theme, "top"))
        col_w = (sp["contentWidth"] - sp["elementGap"]) / 2
        right_x = sp["slideMargin"] + col_w + sp["elementGap"]
        if left_title:
            requests.extend(_build_text_element(
                slide_id, theme, left_title, "heading2",
                sp["slideMargin"], 25, sp["contentWidth"], 40))
        if left_body:
            requests.extend(_build_text_element(
                slide_id, theme, left_body, "body",
                sp["slideMargin"], 75, col_w, 280))
        if right_body:
            requests.extend(_build_text_element(
                slide_id, theme, right_body, "body",
                right_x, 75, col_w, 280))

    elif layout == "closing":
        title = data.get("title", "Thank You")
        subtitle = data.get("subtitle", "")
        requests.append({
            "updatePageProperties": {
                "objectId": slide_id,
                "pageProperties": {"pageBackgroundFill": _solid_fill(colors["primary"])},
                "fields": "pageBackgroundFill",
            }
        })
        requests.extend(_build_accent_bar(slide_id, theme, "bottom"))
        requests.extend(_build_text_element(
            slide_id, theme, title, "title",
            sp["slideMargin"], SLIDE_HEIGHT_PT * 0.3, sp["contentWidth"], 60,
            color_override="#ffffff", align="CENTER"))
        if subtitle:
            requests.extend(_build_text_element(
                slide_id, theme, subtitle, "body",
                sp["slideMargin"], SLIDE_HEIGHT_PT * 0.3 + 65, sp["contentWidth"], 35,
                color_override="#aaaaaa", align="CENTER"))

    elif layout == "content":
        title = data.get("title", "")
        body = data.get("content", data.get("body", ""))
        requests.append({
            "updatePageProperties": {
                "objectId": slide_id,
                "pageProperties": {"pageBackgroundFill": _solid_fill(colors["background"])},
                "fields": "pageBackgroundFill",
            }
        })
        requests.extend(_build_accent_bar(slide_id, theme, "top"))
        if title:
            requests.extend(_build_text_element(
                slide_id, theme, title, "heading1",
                sp["slideMargin"], 20, sp["contentWidth"], 45))
        if body:
            requests.extend(_build_text_element(
                slide_id, theme, body, "body",
                sp["slideMargin"], 75, sp["contentWidth"], 310))

    elif layout == "ad-card":
        # Ad creative card: image + metrics + preview link (1 ad per slide)
        title = data.get("title", "Top Performer")
        ad_name = data.get("ad_name", "")
        campaign = data.get("campaign", "")
        image_url = data.get("image_url", "")
        preview_link = data.get("preview_link", "")
        metrics = data.get("metrics", [])

        # Light background
        requests.append({
            "updatePageProperties": {
                "objectId": slide_id,
                "pageProperties": {"pageBackgroundFill": _solid_fill(colors["background"])},
                "fields": "pageBackgroundFill",
            }
        })
        requests.extend(_build_accent_bar(slide_id, theme, "top"))

        # Title
        requests.extend(_build_text_element(
            slide_id, theme, title, "heading1",
            sp["slideMargin"], 15, sp["contentWidth"], 35, align="START"))

        margin = sp["slideMargin"]
        img_w = 280  # ~45% of slide width
        img_h = 210
        img_y = 60
        right_x = margin + img_w + 20  # gap after image
        right_w = sp["contentWidth"] - img_w - 20

        if image_url:
            # Creative image on the left
            img_id = _gen_id()
            requests.append({
                "createImage": {
                    "objectId": img_id,
                    "url": image_url,
                    "elementProperties": {
                        "pageObjectId": slide_id,
                        "size": {
                            "height": _emu(_emu_from_pt(img_h)),
                            "width": _emu(_emu_from_pt(img_w)),
                        },
                        "transform": {
                            "scaleX": 1, "scaleY": 1,
                            "translateX": _emu_from_pt(margin),
                            "translateY": _emu_from_pt(img_y),
                            "unit": "EMU",
                        },
                    },
                }
            })
        else:
            # No image — shift metrics to full width
            right_x = margin
            right_w = sp["contentWidth"]

        # Ad name (right side, top)
        if ad_name:
            requests.extend(_build_text_element(
                slide_id, theme, ad_name, "heading2",
                right_x, img_y, right_w, 30, align="START"))

        # Campaign name
        if campaign:
            requests.extend(_build_text_element(
                slide_id, theme, campaign, "body",
                right_x, img_y + 32, right_w, 22,
                color_override=colors.get("textSecondary", "#888888"), align="START"))

        # Preview link
        if preview_link:
            requests.extend(_build_text_element(
                slide_id, theme, f"Preview: {preview_link}", "body",
                right_x, img_y + 56, right_w, 20,
                color_override=colors.get("accent", "#2563eb"), align="START"))

        # Metrics cards: 2 rows x 4 cols on the right side below text
        metrics_y = img_y + 85
        n_cols = 4
        n_metrics = min(len(metrics), 8)
        if n_metrics > 0:
            card_gap = 8
            card_w = (right_w - card_gap * (n_cols - 1)) / n_cols
            card_h = 55
            for idx, met in enumerate(metrics[:8]):
                row = idx // n_cols
                col = idx % n_cols
                cx = right_x + col * (card_w + card_gap)
                cy = metrics_y + row * (card_h + card_gap)
                requests.extend(_build_kpi_card(
                    slide_id, theme,
                    met.get("value", "—"), met.get("label", "—"),
                    cx, cy, card_w, card_h))

    else:
        print(f"ERROR: Unknown layout type: {layout}", file=sys.stderr)
        sys.exit(1)

    svc.presentations().batchUpdate(
        presentationId=args.presentation_id, body={"requests": requests}
    ).execute()

    # Auto-clean blank first slide (left behind by create --clean + add-layout)
    try:
        pres = svc.presentations().get(presentationId=args.presentation_id).execute()
        slides = pres.get("slides", [])
        if len(slides) > 1:
            first = slides[0]
            if not first.get("pageElements", []) and first["objectId"] != slide_id:
                svc.presentations().batchUpdate(
                    presentationId=args.presentation_id,
                    body={"requests": [{"deleteObject": {"objectId": first["objectId"]}}]}
                ).execute()
    except Exception:
        pass  # Non-critical cosmetic fix

    print(json.dumps({"slideId": slide_id, "layout": layout}))


def cmd_slides_export(args):
    """Export a presentation as PDF, PPTX, PNG, or SVG."""
    drive = _drive_service()
    mime_map = {
        "pdf": "application/pdf",
        "pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        "png": "image/png",
        "svg": "image/svg+xml",
    }
    mime = mime_map.get(args.format)
    if not mime:
        print(f"ERROR: Unsupported format '{args.format}'. Use: {', '.join(mime_map.keys())}", file=sys.stderr)
        sys.exit(1)

    content = drive.files().export(fileId=args.presentation_id, mimeType=mime).execute()
    out_path = args.output or f"presentation.{args.format}"
    pathlib.Path(out_path).write_bytes(content)
    print(json.dumps({"exported": out_path, "format": args.format, "bytes": len(content)}))


def cmd_slides_batch_update(args):
    """Execute raw batchUpdate requests on a presentation."""
    svc = _slides_service()
    requests_data = _parse_json_arg(args.requests, "requests")
    if not isinstance(requests_data, list):
        requests_data = [requests_data]

    resp = svc.presentations().batchUpdate(
        presentationId=args.presentation_id, body={"requests": requests_data}
    ).execute()
    print(json.dumps({"replies": len(resp.get("replies", []))}))


def _build_rich_presentation(spec, output_path):
    """Build a rich .pptx file from a JSON spec using python-pptx.

    spec format:
    {
      "title": "Presentation Title",
      "theme": {"dark": "1A1A2E", "accent": "E94560", "accent2": "0F3460",
                "lightBg": "F5F5F5", "medGray": "666666"},
      "slideWidth": 13.333,
      "slideHeight": 7.5,
      "slides": [
        {"type": "title_slide", "title": "...", "subtitle": "...", "notes": "..."},
        {"type": "section", "title": "..."},
        {"type": "content", "title": "...", "body": "..."},
        {"type": "bullets", "title": "...", "items": [...]},
        {"type": "kpi_row", "title": "...", "metrics": [{"value":"$42K","label":"Revenue","bg":"0F3460"}]},
        {"type": "table", "title": "...", "headers": [...], "rows": [[...]], "style": {"headerBg":"...","zebra":true}},
        {"type": "two_column", "title": "...", "left": {"heading":"...","body":"..."}, "right": {"heading":"...","body":"..."}},
        {"type": "image", "title": "...", "path": "...", "caption": "..."},
        {"type": "chart", "title": "...", "chartType": "bar", "categories": [...], "series": [{"name":"...","values":[...]}]},
        {"type": "blank", "elements": [{"type":"text","text":"...","x":1,"y":1,"fontSize":14}]},
        {"type": "closing", "title": "...", "subtitle": "..."}
      ]
    }
    """
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.enum.chart import XL_CHART_TYPE
    from pptx.chart.data import CategoryChartData, XyChartData

    prs = Presentation()

    # Slide dimensions (default: widescreen 16:9)
    slide_w = spec.get("slideWidth", 13.333)
    slide_h = spec.get("slideHeight", 7.5)
    prs.slide_width = Inches(slide_w)
    prs.slide_height = Inches(slide_h)

    # Theme colors
    theme = spec.get("theme", {})
    color_dark = theme.get("dark", "1A1A2E")
    color_accent = theme.get("accent", "E94560")
    color_accent2 = theme.get("accent2", "0F3460")
    color_light_bg = theme.get("lightBg", "F5F5F5")
    color_med_gray = theme.get("medGray", "666666")

    def _rgb(hex_str):
        h = hex_str.lstrip("#")
        try:
            return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))
        except (ValueError, IndexError):
            return RGBColor(0, 0, 0)

    def _set_slide_bg(slide, hex_color):
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = _rgb(hex_color)

    def _add_textbox(slide, text, left, top, width, height,
                     font_size=18, bold=False, italic=False,
                     color=None, align=PP_ALIGN.LEFT, font_name="Arial",
                     anchor=MSO_ANCHOR.TOP):
        txBox = slide.shapes.add_textbox(
            Inches(left), Inches(top), Inches(width), Inches(height))
        tf = txBox.text_frame
        tf.word_wrap = True
        tf.auto_size = None
        try:
            tf.paragraphs[0].alignment = align
        except Exception:
            pass
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = text
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.name = font_name
        if color:
            run.font.color.rgb = _rgb(color)
        return txBox

    def _add_accent_bar(slide, position="top", hex_color=None):
        from pptx.util import Inches, Pt
        bar_color = hex_color or color_accent
        bar_h = 0.06
        if position == "top":
            left, top, w, h = 0, 0, slide_w, bar_h
        elif position == "bottom":
            left, top, w, h = 0, slide_h - bar_h, slide_w, bar_h
        else:
            left, top, w, h = 0, 0, slide_w, bar_h
        shape = slide.shapes.add_shape(
            1,  # MSO_SHAPE.RECTANGLE
            Inches(left), Inches(top), Inches(w), Inches(h))
        shape.fill.solid()
        shape.fill.fore_color.rgb = _rgb(bar_color)
        shape.line.fill.background()
        return shape

    def _add_notes(slide, notes_text):
        if notes_text:
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = notes_text

    def _truncate(text, max_chars):
        """Truncate text with ellipsis if too long."""
        if not text or len(text) <= max_chars:
            return text
        return text[:max_chars - 3].rstrip() + "..."

    def _build_kpi_cards(slide, metrics, start_y=2.0, title_text=None):
        margin = 0.8
        content_w = slide_w - 2 * margin
        if title_text:
            _add_textbox(slide, title_text,
                         margin, 0.6, content_w, 0.6,
                         font_size=28, bold=True, color=color_dark)

        if not metrics:
            return
        n = len(metrics)
        max_per_row = 4
        rows_needed = (n + max_per_row - 1) // max_per_row

        gap = 0.2
        padding = 0.15  # inner padding for text
        card_h = 1.5 if rows_needed == 1 else 1.3
        row_gap = 0.25

        for row_idx in range(rows_needed):
            row_start = row_idx * max_per_row
            row_end = min(row_start + max_per_row, n)
            row_count = row_end - row_start
            card_w = (content_w - gap * (row_count - 1)) / row_count
            row_y = start_y + row_idx * (card_h + row_gap)

            # Scale font sizes based on card width
            val_font = max(20, min(32, int(card_w * 10)))
            label_font = max(11, min(14, int(card_w * 4.5)))

            for i, kpi in enumerate(metrics[row_start:row_end]):
                bg = kpi.get("bg", color_accent2)
                val = _truncate(str(kpi.get("value", "--")), 20)
                label = _truncate(str(kpi.get("label", "")), 30)
                card_x = margin + i * (card_w + gap)

                # Card background shape
                shape = slide.shapes.add_shape(
                    1, Inches(card_x), Inches(row_y),
                    Inches(card_w), Inches(card_h))
                shape.fill.solid()
                shape.fill.fore_color.rgb = _rgb(bg)
                shape.line.fill.background()

                # Value text (with padding)
                _add_textbox(slide, val,
                             card_x + padding, row_y + padding,
                             card_w - 2 * padding, card_h * 0.55,
                             font_size=val_font, bold=True, color="FFFFFF",
                             align=PP_ALIGN.CENTER)

                # Label text (with padding)
                _add_textbox(slide, label,
                             card_x + padding, row_y + card_h * 0.6,
                             card_w - 2 * padding, card_h * 0.35,
                             font_size=label_font, color="FFFFFF",
                             align=PP_ALIGN.CENTER)

    def _build_table_slide(slide, headers, rows, style_spec=None, title_text=None):
        margin = 0.8
        content_w = slide_w - 2 * margin
        style_spec = style_spec or {}
        header_bg = style_spec.get("headerBg", color_accent2)
        zebra = style_spec.get("zebra", False)
        zebra_color = style_spec.get("zebraColor", "F0F0F0")

        if title_text:
            _add_textbox(slide, title_text,
                         margin, 0.4, content_w, 0.6,
                         font_size=28, bold=True, color=color_dark)

        all_rows = [headers] + rows if headers else rows
        if not all_rows:
            return
        n_rows = len(all_rows)
        n_cols = max(len(r) for r in all_rows)
        table_top = 1.3 if title_text else 0.6
        table_h = min(n_rows * 0.45, slide_h - table_top - 0.5)

        tbl_shape = slide.shapes.add_table(
            n_rows, n_cols,
            Inches(margin), Inches(table_top),
            Inches(content_w), Inches(table_h))
        tbl = tbl_shape.table

        for r_idx, row_data in enumerate(all_rows):
            for c_idx in range(n_cols):
                cell = tbl.cell(r_idx, c_idx)
                cell_val = row_data[c_idx] if c_idx < len(row_data) else ""
                cell.text = _truncate(str(cell_val), 80)

                for paragraph in cell.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.size = Pt(11)
                        run.font.name = "Arial"

                if r_idx == 0 and headers:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = _rgb(header_bg)
                    for paragraph in cell.text_frame.paragraphs:
                        for run in paragraph.runs:
                            run.font.bold = True
                            run.font.color.rgb = RGBColor(255, 255, 255)
                elif zebra and r_idx % 2 == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = _rgb(zebra_color)

    # Use blank layout for all slides
    blank_layout = prs.slide_layouts[6]  # Blank layout

    for slide_spec in spec.get("slides", []):
        slide_type = slide_spec.get("type", "blank")
        slide = prs.slides.add_slide(blank_layout)
        notes = slide_spec.get("notes", "")
        margin = 0.8
        content_w = slide_w - 2 * margin

        if slide_type == "title_slide":
            _set_slide_bg(slide, color_dark)
            _add_accent_bar(slide, "top")
            title = slide_spec.get("title", "Untitled Presentation")
            subtitle = slide_spec.get("subtitle", "")
            _add_textbox(slide, title,
                         margin, slide_h * 0.3, content_w, 1.0,
                         font_size=40, bold=True, color="FFFFFF",
                         align=PP_ALIGN.CENTER)
            if subtitle:
                _add_textbox(slide, subtitle,
                             margin, slide_h * 0.3 + 1.1, content_w, 0.6,
                             font_size=20, color="CCCCCC",
                             align=PP_ALIGN.CENTER)

        elif slide_type == "section":
            _set_slide_bg(slide, color_accent2)
            _add_accent_bar(slide, "bottom")
            title = slide_spec.get("title", "Section")
            _add_textbox(slide, title,
                         margin, slide_h * 0.35, content_w, 1.0,
                         font_size=36, bold=True, color="FFFFFF",
                         align=PP_ALIGN.CENTER)

        elif slide_type == "content":
            _set_slide_bg(slide, color_light_bg)
            _add_accent_bar(slide, "top")
            title = slide_spec.get("title", "")
            body = _truncate(slide_spec.get("body", slide_spec.get("content", "")), 600)
            if title:
                _add_textbox(slide, title,
                             margin, 0.4, content_w, 0.6,
                             font_size=28, bold=True, color=color_dark)
            if body:
                _add_textbox(slide, body,
                             margin, 1.2, content_w, slide_h - 1.8,
                             font_size=16, color="333333")

        elif slide_type == "bullets":
            _set_slide_bg(slide, color_light_bg)
            _add_accent_bar(slide, "top")
            title = slide_spec.get("title", "")
            items = slide_spec.get("items", [])
            if title:
                _add_textbox(slide, title,
                             margin, 0.4, content_w, 0.6,
                             font_size=28, bold=True, color=color_dark)
            if items:
                # Cap at 10 bullets, indicate truncation
                display_items = items[:10]
                if len(items) > 10:
                    display_items.append(f"... and {len(items) - 10} more")
                txBox = slide.shapes.add_textbox(
                    Inches(margin), Inches(1.3),
                    Inches(content_w), Inches(slide_h - 1.8))
                tf = txBox.text_frame
                tf.word_wrap = True
                # Scale font if many bullets
                bullet_font = 16 if len(display_items) <= 7 else max(12, 16 - (len(display_items) - 7))
                for idx, item in enumerate(display_items):
                    if idx == 0:
                        p = tf.paragraphs[0]
                    else:
                        p = tf.add_paragraph()
                    p.space_after = Pt(6 if len(display_items) > 7 else 8)
                    p.level = 0
                    run = p.add_run()
                    run.text = "\u2022  " + _truncate(str(item), 120)
                    run.font.size = Pt(bullet_font)
                    run.font.name = "Arial"
                    run.font.color.rgb = _rgb("333333")

        elif slide_type == "kpi_row":
            _set_slide_bg(slide, color_light_bg)
            _add_accent_bar(slide, "top")
            title = slide_spec.get("title", "Key Metrics")
            metrics = slide_spec.get("metrics", slide_spec.get("items", []))
            _build_kpi_cards(slide, metrics, start_y=2.2, title_text=title)

        elif slide_type == "table":
            _set_slide_bg(slide, color_light_bg)
            _add_accent_bar(slide, "top")
            title = slide_spec.get("title", "")
            headers = slide_spec.get("headers", [])
            rows = slide_spec.get("rows", [])
            style = slide_spec.get("style", {})
            _build_table_slide(slide, headers, rows, style, title)

        elif slide_type == "two_column":
            _set_slide_bg(slide, color_light_bg)
            _add_accent_bar(slide, "top")
            title = slide_spec.get("title", "")
            left = slide_spec.get("left", {})
            right = slide_spec.get("right", {})
            col_w = (content_w - 0.4) / 2

            if title:
                _add_textbox(slide, title,
                             margin, 0.4, content_w, 0.6,
                             font_size=28, bold=True, color=color_dark)

            # Left column
            left_heading = left.get("heading", "")
            left_body = _truncate(left.get("body", ""), 400)
            y_start = 1.3
            if left_heading:
                _add_textbox(slide, left_heading,
                             margin, y_start, col_w, 0.5,
                             font_size=20, bold=True, color=color_accent2)
                y_start += 0.55
            if left_body:
                _add_textbox(slide, left_body,
                             margin, y_start, col_w, slide_h - y_start - 0.5,
                             font_size=14, color="333333")

            # Right column
            right_heading = right.get("heading", "")
            right_body = _truncate(right.get("body", ""), 400)
            right_x = margin + col_w + 0.4
            y_start = 1.3
            if right_heading:
                _add_textbox(slide, right_heading,
                             right_x, y_start, col_w, 0.5,
                             font_size=20, bold=True, color=color_accent2)
                y_start += 0.55
            if right_body:
                _add_textbox(slide, right_body,
                             right_x, y_start, col_w, slide_h - y_start - 0.5,
                             font_size=14, color="333333")

        elif slide_type == "image":
            _set_slide_bg(slide, color_light_bg)
            _add_accent_bar(slide, "top")
            title = slide_spec.get("title", "")
            img_path = slide_spec.get("path", "")
            caption = slide_spec.get("caption", "")

            if title:
                _add_textbox(slide, title,
                             margin, 0.4, content_w, 0.6,
                             font_size=28, bold=True, color=color_dark)

            if img_path and os.path.exists(img_path):
                img_top = 1.3 if title else 0.6
                max_w = content_w
                max_h = slide_h - img_top - (0.8 if caption else 0.4)
                slide.shapes.add_picture(
                    img_path,
                    Inches(margin), Inches(img_top),
                    Inches(max_w), Inches(max_h))

            if caption:
                _add_textbox(slide, caption,
                             margin, slide_h - 0.7, content_w, 0.4,
                             font_size=11, italic=True, color=color_med_gray,
                             align=PP_ALIGN.CENTER)

        elif slide_type == "chart":
            _set_slide_bg(slide, color_light_bg)
            _add_accent_bar(slide, "top")
            title = slide_spec.get("title", "")
            chart_type_str = slide_spec.get("chartType", "bar").lower()
            categories = slide_spec.get("categories", [])
            series_list = slide_spec.get("series", [])

            if title:
                _add_textbox(slide, title,
                             margin, 0.4, content_w, 0.6,
                             font_size=28, bold=True, color=color_dark)

            # Guard: need at least one series with data
            valid_series = [s for s in series_list if s.get("values")]
            if not valid_series:
                _add_textbox(slide, "(No chart data provided)",
                             margin, 3.0, content_w, 1.0,
                             font_size=16, italic=True, color=color_med_gray,
                             align=PP_ALIGN.CENTER)
            else:
                chart_type_map = {
                    "bar": XL_CHART_TYPE.COLUMN_CLUSTERED,
                    "column": XL_CHART_TYPE.COLUMN_CLUSTERED,
                    "line": XL_CHART_TYPE.LINE,
                    "pie": XL_CHART_TYPE.PIE,
                    "area": XL_CHART_TYPE.AREA,
                    "scatter": XL_CHART_TYPE.XY_SCATTER,
                    "doughnut": XL_CHART_TYPE.DOUGHNUT,
                }
                xl_type = chart_type_map.get(chart_type_str, XL_CHART_TYPE.COLUMN_CLUSTERED)

                # Scatter charts require XyChartData; everything else uses CategoryChartData
                if chart_type_str == "scatter":
                    chart_data = XyChartData()
                    for s in valid_series:
                        xy_series = chart_data.add_series(s.get("name", "Series"))
                        x_vals = categories if categories else list(range(1, len(s["values"]) + 1))
                        for x, y in zip(x_vals, s["values"]):
                            xy_series.add_data_point(x, y)
                else:
                    chart_data = CategoryChartData()
                    if not categories:
                        categories = [str(i + 1) for i in range(len(valid_series[0]["values"]))]
                    chart_data.categories = categories
                    for s in valid_series:
                        chart_data.add_series(s.get("name", "Series"), s.get("values", []))

                chart_top = 1.3 if title else 0.5
                chart_shape = slide.shapes.add_chart(
                    xl_type, Inches(margin), Inches(chart_top),
                    Inches(content_w), Inches(slide_h - chart_top - 0.5),
                    chart_data)
                chart_obj = chart_shape.chart
                chart_obj.has_legend = len(valid_series) > 1

        elif slide_type == "blank":
            bg_color = slide_spec.get("background", color_light_bg)
            _set_slide_bg(slide, bg_color)
            for elem in slide_spec.get("elements", []):
                etype = elem.get("type", "text")
                if etype == "text":
                    _add_textbox(
                        slide, elem.get("text", ""),
                        elem.get("x", 0.5), elem.get("y", 0.5),
                        elem.get("width", content_w), elem.get("height", 1.0),
                        font_size=elem.get("fontSize", 14),
                        bold=elem.get("bold", False),
                        italic=elem.get("italic", False),
                        color=elem.get("color", "333333"),
                        align=PP_ALIGN.CENTER if elem.get("align") == "center"
                              else PP_ALIGN.RIGHT if elem.get("align") == "right"
                              else PP_ALIGN.LEFT)
                elif etype == "image" and elem.get("path") and os.path.exists(elem["path"]):
                    slide.shapes.add_picture(
                        elem["path"],
                        Inches(elem.get("x", 0.5)), Inches(elem.get("y", 0.5)),
                        Inches(elem.get("width", 4)), Inches(elem.get("height", 3)))
                elif etype == "shape":
                    shape = slide.shapes.add_shape(
                        1,
                        Inches(elem.get("x", 0)), Inches(elem.get("y", 0)),
                        Inches(elem.get("width", 1)), Inches(elem.get("height", 1)))
                    shape.fill.solid()
                    shape.fill.fore_color.rgb = _rgb(elem.get("color", color_accent))
                    shape.line.fill.background()

        elif slide_type == "closing":
            _set_slide_bg(slide, color_dark)
            _add_accent_bar(slide, "bottom")
            title = slide_spec.get("title", "Thank You")
            subtitle = slide_spec.get("subtitle", "")
            _add_textbox(slide, title,
                         margin, slide_h * 0.3, content_w, 1.0,
                         font_size=40, bold=True, color="FFFFFF",
                         align=PP_ALIGN.CENTER)
            if subtitle:
                _add_textbox(slide, subtitle,
                             margin, slide_h * 0.3 + 1.1, content_w, 0.6,
                             font_size=20, color="AAAAAA",
                             align=PP_ALIGN.CENTER)

        _add_notes(slide, notes)

    prs.save(output_path)
    return output_path


def _qa_presentation(pptx_path, spec):
    """Vision QA self-correction loop: export to images and check for layout issues."""
    import subprocess

    report = []
    pdf_path = str(pptx_path).rsplit(".", 1)[0] + ".pdf"
    images_dir = str(pptx_path).rsplit(".", 1)[0] + "_qa_images"

    # Step 1: Convert .pptx to PDF
    try:
        for cmd in [
            ["libreoffice", "--headless", "--convert-to", "pdf", "--outdir",
             os.path.dirname(pdf_path), str(pptx_path)],
            ["unoconv", "-f", "pdf", "-o", pdf_path, str(pptx_path)],
        ]:
            try:
                subprocess.run(cmd, capture_output=True, timeout=60, check=True)
                if os.path.exists(pdf_path):
                    break
            except (FileNotFoundError, subprocess.CalledProcessError):
                continue
        if not os.path.exists(pdf_path):
            return {"skipped": True, "reason": "No PDF converter available (install LibreOffice or unoconv)"}
    except Exception as e:
        return {"skipped": True, "reason": f"PDF conversion failed: {e}"}

    # Step 2: Convert PDF pages to PNG images
    os.makedirs(images_dir, exist_ok=True)
    slide_images = []
    try:
        try:
            import fitz  # pymupdf
            doc = fitz.open(pdf_path)
            for page_num in range(len(doc)):
                pix = doc[page_num].get_pixmap(dpi=150)
                img_path = os.path.join(images_dir, f"slide_{page_num}.png")
                pix.save(img_path)
                slide_images.append(img_path)
            doc.close()
        except ImportError:
            from pdf2image import convert_from_path
            images = convert_from_path(pdf_path, dpi=150)
            for i, img in enumerate(images):
                img_path = os.path.join(images_dir, f"slide_{i}.png")
                img.save(img_path, "PNG")
                slide_images.append(img_path)
    except Exception as e:
        return {"skipped": True, "reason": f"PDF-to-image conversion failed: {e}"}

    if not slide_images:
        return {"skipped": True, "reason": "No slide images generated"}

    # Step 3: Vision QA via Anthropic API
    try:
        import anthropic
        api_key = os.environ.get("ANTHROPIC_API_KEY")
        if not api_key:
            return {"skipped": True, "reason": "ANTHROPIC_API_KEY not set"}
        client = anthropic.Anthropic(api_key=api_key)
    except ImportError:
        return {"skipped": True, "reason": "anthropic package not installed"}

    import base64
    qa_prompt = (
        "Analyze this presentation slide for layout issues: overlapping text, "
        "text overflow, empty charts, broken formatting. List specific issues "
        "found or say PASS if the slide looks correct."
    )

    for idx, img_path in enumerate(slide_images):
        try:
            with open(img_path, "rb") as f:
                img_b64 = base64.b64encode(f.read()).decode("utf-8")
            resp = client.messages.create(
                model="claude-haiku-4-5-20251001",
                max_tokens=300,
                messages=[{
                    "role": "user",
                    "content": [
                        {"type": "image", "source": {"type": "base64", "media_type": "image/png", "data": img_b64}},
                        {"type": "text", "text": qa_prompt},
                    ],
                }],
            )
            result_text = resp.content[0].text.strip()
            status = "PASS" if result_text.upper().startswith("PASS") else "FAIL"
            issues = [] if status == "PASS" else [result_text]
            report.append({"slide_index": idx, "status": status, "issues": issues})
        except Exception as e:
            report.append({"slide_index": idx, "status": "ERROR", "issues": [str(e)]})

    # Cleanup temp files
    try:
        os.remove(pdf_path)
        import shutil
        shutil.rmtree(images_dir, ignore_errors=True)
    except Exception:
        pass

    return {"slides": report, "has_issues": any(s["status"] != "PASS" for s in report)}


def _upload_pptx_as_slides(pptx_path, title, folder_id=None):
    """Upload .pptx to Drive, converting to Google Slides. Returns file metadata."""
    from googleapiclient.http import MediaFileUpload
    drive = _drive_service()

    metadata = {
        "name": title,
        "mimeType": "application/vnd.google-apps.presentation",
    }
    if folder_id:
        metadata["parents"] = [folder_id]

    media = MediaFileUpload(
        str(pptx_path),
        mimetype="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        resumable=True,
    )
    uploaded = drive.files().create(
        body=metadata,
        media_body=media,
        fields="id, name, webViewLink",
    ).execute()
    return uploaded


def cmd_slides_build_rich(args):
    """Build a rich presentation from JSON spec, upload to Drive as Google Slides."""
    spec = _read_spec(args)
    title = args.title or spec.get("title", "Untitled Presentation")

    # Build .pptx locally
    if args.output:
        output_path = pathlib.Path(args.output)
    else:
        output_path = pathlib.Path(tempfile.mkdtemp()) / f"{title}.pptx"

    _build_rich_presentation(spec, str(output_path))

    # Optional QA check
    qa_result = None
    if getattr(args, "qa", False):
        qa_result = _qa_presentation(str(output_path), spec)
        if qa_result and not qa_result.get("skipped"):
            report_path = str(output_path) + ".qa-report.json"
            with open(report_path, "w") as f:
                json.dump(qa_result, f, indent=2)
            if qa_result.get("has_issues"):
                print(f"QA WARNING: Issues found -- see {report_path}", file=sys.stderr)
        elif qa_result and qa_result.get("skipped"):
            print(f"QA skipped: {qa_result.get('reason', 'unknown')}", file=sys.stderr)

    if args.output:
        # Local-only mode -- skip upload
        result = {"localPath": str(output_path), "title": title}
        if qa_result:
            result["qa"] = qa_result
        print(json.dumps(result))
        return

    # Upload to Drive as Google Slides
    uploaded = _upload_pptx_as_slides(str(output_path), title, folder_id=args.folder)
    pres_id = uploaded["id"]
    _auto_share_org(pres_id)
    url = uploaded.get("webViewLink", f"https://docs.google.com/presentation/d/{pres_id}/edit")

    result = {"presentationId": pres_id, "url": url, "title": title, "orgShared": True}

    if args.link_share:
        result["linkShare"] = _link_share_file(pres_id)

    if args.share:
        drive = _drive_service()
        drive.permissions().create(
            fileId=pres_id,
            body={"type": "user", "role": "reader", "emailAddress": args.share},
            sendNotificationEmail=True,
            fields="id",
        ).execute()
        result["sharedWith"] = args.share

    # Clean up temp file
    output_path.unlink(missing_ok=True)

    print(json.dumps(result, indent=2))


def cmd_slides_build_rich_batch(args):
    """Build multiple rich presentations in parallel from JSON spec array."""
    from concurrent.futures import ThreadPoolExecutor, as_completed

    specs = _read_spec(args)
    if not isinstance(specs, list):
        print("ERROR: Batch spec must be a JSON array of specs.", file=sys.stderr)
        sys.exit(1)

    max_workers = args.parallel or 4
    results = {"total": len(specs), "succeeded": 0, "failed": 0, "presentations": []}

    def build_one(spec):
        title = spec.get("title", "Untitled Presentation")
        tmp_path = pathlib.Path(tempfile.mkdtemp()) / f"{title}.pptx"
        try:
            _build_rich_presentation(spec, str(tmp_path))
            uploaded = _upload_pptx_as_slides(str(tmp_path), title, folder_id=args.folder)
            pres_id = uploaded["id"]
            _auto_share_org(pres_id)
            url = uploaded.get("webViewLink", f"https://docs.google.com/presentation/d/{pres_id}/edit")
            entry = {"presentationId": pres_id, "url": url, "title": title, "orgShared": True}
            if args.link_share:
                entry["linkShare"] = _link_share_file(pres_id)
            return entry
        except Exception as e:
            return {"title": title, "error": str(e)}
        finally:
            tmp_path.unlink(missing_ok=True)

    with ThreadPoolExecutor(max_workers=max_workers) as pool:
        futures = {pool.submit(build_one, s): s for s in specs}
        for future in as_completed(futures):
            entry = future.result()
            if "error" in entry:
                results["failed"] += 1
            else:
                results["succeeded"] += 1
            results["presentations"].append(entry)

    print(json.dumps(results, indent=2))


# ============================================================
# DOCS commands
# ============================================================

def _doc_end_index(svc, doc_id):
    """Get the end index of a document's body content."""
    doc = svc.documents().get(documentId=doc_id).execute()
    body = doc.get("body", {})
    content = body.get("content", [])
    if content:
        return content[-1].get("endIndex", 1) - 1
    return 1


def _doc_heading_index(svc, doc_id, heading_text):
    """Find the end index of a heading's content block (insert point after it).

    Searches for a heading whose text contains heading_text (case-insensitive).
    Returns the endIndex of the paragraph element containing the heading,
    which is the right insertion point for content after that heading.
    """
    doc = svc.documents().get(documentId=doc_id).execute()
    for elem in doc.get("body", {}).get("content", []):
        para = elem.get("paragraph", {})
        named_style = para.get("paragraphStyle", {}).get("namedStyleType", "")
        if named_style.startswith("HEADING") or named_style in ("TITLE", "SUBTITLE"):
            text = "".join(
                run.get("textRun", {}).get("content", "")
                for run in para.get("elements", [])
            ).strip()
            if heading_text.lower() in text.lower():
                return elem.get("endIndex", 1) - 1
    return None


def _resolve_insert_index(svc, doc_id, args):
    """Determine insertion index from --index or --after-heading, falling back to end.

    Returns (index: int, source: str) where source describes where we're inserting.
    """
    if hasattr(args, "index") and args.index is not None:
        return args.index, f"index:{args.index}"

    if hasattr(args, "after_heading") and args.after_heading:
        idx = _doc_heading_index(svc, doc_id, args.after_heading)
        if idx is not None:
            return idx, f"after-heading:'{args.after_heading}'"
        print(f"WARNING: Heading '{args.after_heading}' not found — appending to end", file=sys.stderr)

    return _doc_end_index(svc, doc_id), "end"


def cmd_docs_create(args):
    """Create a new Google Doc."""
    svc = _docs_service()
    doc = svc.documents().create(body={"title": args.title}).execute()
    doc_id = doc["documentId"]
    _auto_share_org(doc_id)
    url = f"https://docs.google.com/document/d/{doc_id}/edit"
    print(json.dumps({"documentId": doc_id, "url": url, "orgShared": True}))


def cmd_docs_get(args):
    """Get document content and structure."""
    svc = _docs_service()
    doc = svc.documents().get(documentId=args.document_id).execute()
    summary = {
        "documentId": doc["documentId"],
        "title": doc.get("title", ""),
        "url": f"https://docs.google.com/document/d/{doc['documentId']}/edit",
        "endIndex": doc.get("body", {}).get("content", [{}])[-1].get("endIndex", 1),
    }
    # Extract headings for outline
    headings = []
    for elem in doc.get("body", {}).get("content", []):
        para = elem.get("paragraph", {})
        named_style = para.get("paragraphStyle", {}).get("namedStyleType", "")
        if named_style.startswith("HEADING"):
            text = "".join(
                run.get("textRun", {}).get("content", "")
                for run in para.get("elements", [])
            ).strip()
            headings.append({"style": named_style, "text": text, "index": elem.get("startIndex")})
    summary["headings"] = headings
    print(json.dumps(summary, indent=2))


def cmd_docs_add_text(args):
    """Add styled text to a document."""
    # Gate: reject markdown content — use structured commands instead
    if _detect_markdown(args.text):
        print("ERROR: Markdown detected in text content. Do NOT pass markdown to add-text.", file=sys.stderr)
        print("Use the structured commands instead:", file=sys.stderr)
        print("  - add-text with --style heading1/heading2/heading3 for headings", file=sys.stderr)
        print("  - add-text with --text-style '{\"bold\":true}' for bold text", file=sys.stderr)
        print("  - add-list for bullet/numbered lists", file=sys.stderr)
        print("  - add-table for data tables", file=sys.stderr)
        print("  - add-divider for horizontal rules", file=sys.stderr)
        print("Each content block should be a SEPARATE add-text/add-list/add-table call.", file=sys.stderr)
        sys.exit(1)

    svc = _docs_service()
    idx, insert_at = _resolve_insert_index(svc, args.document_id, args)

    text = args.text + "\n"
    requests = [
        {"insertText": {"location": {"index": idx}, "text": text}},
    ]

    # Map style names to namedStyleType
    style_map = {
        "heading1": "HEADING_1",
        "heading2": "HEADING_2",
        "heading3": "HEADING_3",
        "heading4": "HEADING_4",
        "heading5": "HEADING_5",
        "heading6": "HEADING_6",
        "title": "TITLE",
        "subtitle": "SUBTITLE",
        "normal": "NORMAL_TEXT",
    }

    named_style = style_map.get(args.style, "NORMAL_TEXT") if args.style else "NORMAL_TEXT"

    requests.append({
        "updateParagraphStyle": {
            "range": {"startIndex": idx, "endIndex": idx + len(text)},
            "paragraphStyle": {"namedStyleType": named_style},
            "fields": "namedStyleType",
        }
    })

    # Apply custom text styling if provided via --text-style
    if args.text_style:
        ts = _parse_json_arg(args.text_style, "text-style")
        style_obj = {}
        fields = []
        if "bold" in ts:
            style_obj["bold"] = ts["bold"]
            fields.append("bold")
        if "italic" in ts:
            style_obj["italic"] = ts["italic"]
            fields.append("italic")
        if "color" in ts:
            style_obj["foregroundColor"] = _rgb_color(ts["color"])
            fields.append("foregroundColor")
        if "fontSize" in ts:
            style_obj["fontSize"] = _pt(ts["fontSize"])
            fields.append("fontSize")
        if "fontFamily" in ts:
            style_obj["weightedFontFamily"] = {"fontFamily": ts["fontFamily"]}
            fields.append("weightedFontFamily")
        if fields:
            requests.append({
                "updateTextStyle": {
                    "range": {"startIndex": idx, "endIndex": idx + len(text) - 1},
                    "textStyle": style_obj,
                    "fields": ",".join(fields),
                }
            })

    svc.documents().batchUpdate(
        documentId=args.document_id, body={"requests": requests}
    ).execute()
    print(json.dumps({"inserted": len(text) - 1, "style": named_style, "at": insert_at}))


def _detect_markdown(text: str) -> bool:
    """Detect if text contains markdown formatting."""
    md_patterns = [
        r"^#{1,6}\s+",          # headings
        r"\*\*[^*]+\*\*",       # bold
        r"^\s*[-*+]\s+",        # bullet lists
        r"^\s*\d+\.\s+",        # numbered lists
        r"```",                  # code blocks
        r"\[.+\]\(.+\)",        # links
    ]
    lines = text.split("\n") if "\n" in text else [text]
    matches = 0
    for line in lines:
        for pattern in md_patterns:
            if re.search(pattern, line):
                matches += 1
                break
    # If >25% of non-empty lines look like markdown, reject
    non_empty = sum(1 for l in lines if l.strip())
    return non_empty > 0 and matches / non_empty > 0.25


def cmd_docs_add_table(args):
    """Add a styled table to a document."""
    svc = _docs_service()
    data = _parse_json_arg(args.data, "data")
    style = _parse_json_arg(args.style, "style") if args.style else {}

    if not data or not isinstance(data, list):
        print("ERROR: --data must be a 2D JSON array.", file=sys.stderr)
        sys.exit(1)

    rows = len(data)
    cols = len(data[0]) if data else 0
    idx, _ = _resolve_insert_index(svc, args.document_id, args)

    header_color = style.get("headerColor", "#1a73e8")

    # Insert table
    requests = [
        {"insertTable": {"location": {"index": idx}, "rows": rows, "columns": cols}},
    ]
    svc.documents().batchUpdate(
        documentId=args.document_id, body={"requests": requests}
    ).execute()

    # Re-read doc to find table structure
    doc = svc.documents().get(documentId=args.document_id).execute()
    tables = []
    for elem in doc.get("body", {}).get("content", []):
        if "table" in elem:
            tables.append(elem)

    if not tables:
        print("ERROR: Table not found after insertion.", file=sys.stderr)
        sys.exit(1)

    table = tables[-1]["table"]
    table_start = tables[-1]["startIndex"]

    # Fill cells with data (insert backwards to maintain indices)
    fill_requests = []
    for r_idx in range(rows - 1, -1, -1):
        row_cells = table["tableRows"][r_idx]["tableCells"]
        for c_idx in range(cols - 1, -1, -1):
            cell = row_cells[c_idx]
            cell_start = cell["content"][0]["startIndex"]
            cell_val = str(data[r_idx][c_idx]) if c_idx < len(data[r_idx]) else ""
            if cell_val:
                fill_requests.append({
                    "insertText": {"location": {"index": cell_start}, "text": cell_val}
                })

    if fill_requests:
        svc.documents().batchUpdate(
            documentId=args.document_id, body={"requests": fill_requests}
        ).execute()

    # Re-read doc again to style header row
    doc = svc.documents().get(documentId=args.document_id).execute()
    tables = [e for e in doc.get("body", {}).get("content", []) if "table" in e]
    if not tables:
        print(json.dumps({"tableInserted": True, "rows": rows, "cols": cols}))
        return

    table = tables[-1]["table"]
    style_requests = []

    # Bold header text + background color
    header_cells = table["tableRows"][0]["tableCells"]
    for c_idx, cell in enumerate(header_cells):
        cell_start = cell["content"][0]["startIndex"]
        cell_end = cell["content"][-1]["endIndex"]
        # Bold the header text
        style_requests.append({
            "updateTextStyle": {
                "range": {"startIndex": cell_start, "endIndex": cell_end - 1},
                "textStyle": {"bold": True},
                "fields": "bold",
            }
        })

    # Header row background
    rgb = _hex_to_rgb(header_color)
    style_requests.append({
        "updateTableCellStyle": {
            "tableRange": {
                "tableCellLocation": {
                    "tableStartLocation": {"index": tables[-1]["startIndex"]},
                    "rowIndex": 0,
                    "columnIndex": 0,
                },
                "rowSpan": 1,
                "columnSpan": cols,
            },
            "tableCellStyle": {
                "backgroundColor": {
                    "color": {"rgbColor": rgb}
                },
            },
            "fields": "backgroundColor",
        }
    })

    # Pin header row
    style_requests.append({
        "pinTableHeaderRows": {
            "tableStartLocation": {"index": tables[-1]["startIndex"]},
            "pinnedHeaderRowsCount": 1,
        }
    })

    if style_requests:
        svc.documents().batchUpdate(
            documentId=args.document_id, body={"requests": style_requests}
        ).execute()

    print(json.dumps({"tableInserted": True, "rows": rows, "cols": cols, "headerStyled": True}))


def cmd_docs_add_image(args):
    """Insert an inline image into a document."""
    svc = _docs_service()
    idx, _ = _resolve_insert_index(svc, args.document_id, args)

    request = {
        "insertInlineImage": {
            "location": {"index": idx},
            "uri": args.url,
        }
    }
    if args.width or args.height:
        request["insertInlineImage"]["objectSize"] = {}
        if args.width:
            request["insertInlineImage"]["objectSize"]["width"] = _pt(args.width)
        if args.height:
            request["insertInlineImage"]["objectSize"]["height"] = _pt(args.height)

    svc.documents().batchUpdate(
        documentId=args.document_id, body={"requests": [request]}
    ).execute()
    print(json.dumps({"imageInserted": True}))


def cmd_docs_add_section(args):
    """Add a section break + heading."""
    svc = _docs_service()
    idx = _doc_end_index(svc, args.document_id)

    requests = [
        {"insertSectionBreak": {"location": {"index": idx}, "sectionType": "NEXT_PAGE"}},
    ]
    svc.documents().batchUpdate(
        documentId=args.document_id, body={"requests": requests}
    ).execute()

    # Now add the heading text in the new section
    idx = _doc_end_index(svc, args.document_id)
    text = args.title + "\n"
    requests = [
        {"insertText": {"location": {"index": idx}, "text": text}},
        {
            "updateParagraphStyle": {
                "range": {"startIndex": idx, "endIndex": idx + len(text)},
                "paragraphStyle": {"namedStyleType": "HEADING_1"},
                "fields": "namedStyleType",
            }
        },
    ]
    svc.documents().batchUpdate(
        documentId=args.document_id, body={"requests": requests}
    ).execute()
    print(json.dumps({"sectionAdded": True, "title": args.title}))


def cmd_docs_add_list(args):
    """Add a bulleted or numbered list."""
    svc = _docs_service()
    items = _parse_json_arg(args.items, "items")
    if not isinstance(items, list):
        print("ERROR: --items must be a JSON array of strings.", file=sys.stderr)
        sys.exit(1)

    idx, insert_at = _resolve_insert_index(svc, args.document_id, args)
    text = "\n".join(items) + "\n"

    requests = [
        {"insertText": {"location": {"index": idx}, "text": text}},
        {
            "createParagraphBullets": {
                "range": {"startIndex": idx, "endIndex": idx + len(text)},
                "bulletPreset": "BULLET_DISC_CIRCLE_SQUARE" if args.list_type == "bullet"
                else "NUMBERED_DECIMAL_ALPHA_ROMAN",
            }
        },
    ]
    svc.documents().batchUpdate(
        documentId=args.document_id, body={"requests": requests}
    ).execute()
    print(json.dumps({"listInserted": True, "items": len(items), "type": args.list_type, "at": insert_at}))


def cmd_docs_add_page_break(args):
    """Insert a page break."""
    svc = _docs_service()
    idx = _doc_end_index(svc, args.document_id)
    svc.documents().batchUpdate(
        documentId=args.document_id,
        body={"requests": [{"insertPageBreak": {"location": {"index": idx}}}]},
    ).execute()
    print(json.dumps({"pageBreakInserted": True}))


def cmd_docs_add_divider(args):
    """Add a visual divider (paragraph with bottom border)."""
    svc = _docs_service()
    idx = _doc_end_index(svc, args.document_id)

    text = " \n"
    requests = [
        {"insertText": {"location": {"index": idx}, "text": text}},
        {
            "updateParagraphStyle": {
                "range": {"startIndex": idx, "endIndex": idx + len(text)},
                "paragraphStyle": {
                    "borderBottom": {
                        "color": {"color": {"rgbColor": _hex_to_rgb("#dadce0")}},
                        "width": _pt(1),
                        "padding": _pt(6),
                        "dashStyle": "SOLID",
                    },
                    "spaceBelow": _pt(12),
                },
                "fields": "borderBottom,spaceBelow",
            }
        },
    ]
    svc.documents().batchUpdate(
        documentId=args.document_id, body={"requests": requests}
    ).execute()
    print(json.dumps({"dividerInserted": True}))


def cmd_docs_add_header(args):
    """Add a page header."""
    svc = _docs_service()
    requests = [{"createHeader": {"type": "DEFAULT", "sectionBreakLocation": {"index": 0}}}]
    resp = svc.documents().batchUpdate(
        documentId=args.document_id, body={"requests": requests}
    ).execute()
    header_id = resp["replies"][0]["createHeader"]["headerId"]

    # Insert text into the header
    svc.documents().batchUpdate(
        documentId=args.document_id,
        body={"requests": [{"insertText": {"location": {"segmentId": header_id, "index": 0}, "text": args.text}}]},
    ).execute()
    print(json.dumps({"headerCreated": True, "headerId": header_id}))


def cmd_docs_add_footer(args):
    """Add a page footer."""
    svc = _docs_service()
    requests = [{"createFooter": {"type": "DEFAULT", "sectionBreakLocation": {"index": 0}}}]
    resp = svc.documents().batchUpdate(
        documentId=args.document_id, body={"requests": requests}
    ).execute()
    footer_id = resp["replies"][0]["createFooter"]["footerId"]

    svc.documents().batchUpdate(
        documentId=args.document_id,
        body={"requests": [{"insertText": {"location": {"segmentId": footer_id, "index": 0}, "text": args.text}}]},
    ).execute()
    print(json.dumps({"footerCreated": True, "footerId": footer_id}))


def cmd_docs_export(args):
    """Export a document as PDF, DOCX, etc."""
    drive = _drive_service()
    mime_map = {
        "pdf": "application/pdf",
        "docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        "md": "text/markdown",
        "html": "text/html",
        "epub": "application/epub+zip",
    }
    mime = mime_map.get(args.format)
    if not mime:
        print(f"ERROR: Unsupported format '{args.format}'. Use: {', '.join(mime_map.keys())}", file=sys.stderr)
        sys.exit(1)

    content = drive.files().export(fileId=args.document_id, mimeType=mime).execute()
    out_path = args.output or f"document.{args.format}"
    pathlib.Path(out_path).write_bytes(content)
    print(json.dumps({"exported": out_path, "format": args.format, "bytes": len(content)}))


# ============================================================
# Rich DOCX pipeline (python-docx)
# ============================================================

def _build_rich_document(spec, output_path):
    """Build a rich .docx file from a JSON spec using python-docx.

    spec format:
    {
      "theme": {"dark": "1A1A2E", "accent": "E94560", "accent2": "0F3460",
                "lightBg": "F5F5F5", "medGray": "666666"},
      "header": "Header text",
      "footer": "Footer text",
      "sections": [{
        "children": [
          {"type": "heading", "level": 1, "text": "...", "color": "1A1A2E"},
          {"type": "subtitle", "text": "..."},
          {"type": "paragraph", "text": "...", "runs": [{"text": "...", "bold": true}]},
          {"type": "accentBar", "color": "E94560"},
          {"type": "divider"},
          {"type": "kpiRow", "items": [{"value": "$42K", "label": "Revenue", "bg": "0F3460"}]},
          {"type": "callout", "title": "...", "body": "...", "bg": "E8F5E9"},
          {"type": "table", "headers": [...], "rows": [[...]], "style": {"headerBg": "0F3460", "zebra": true}},
          {"type": "richTable", "headers": [...], "rows": [[{"text":"...","bold":true}]]},
          {"type": "bullets", "items": [...]},
          {"type": "numberedList", "items": [...]},
          {"type": "personaCard", "name": "...", "role": "...", "details": [...]},
          {"type": "conceptHeader", "title": "...", "subtitle": "..."},
          {"type": "tag", "text": "...", "color": "E94560"},
          {"type": "pageBreak"},
          {"type": "spacer"},
          {"type": "image", "path": "...", "width": 5.0}
        ]
      }]
    }
    """
    from docx import Document
    from docx.shared import Inches, Pt, RGBColor, Cm, Emu
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.enum.table import WD_TABLE_ALIGNMENT
    from docx.oxml.ns import qn, nsdecls
    from docx.oxml import parse_xml

    doc = Document()

    # Set default styles
    style = doc.styles["Normal"]
    style.font.name = "Arial"
    style.font.size = Pt(11)
    style.paragraph_format.space_after = Pt(6)

    # Page setup: US Letter, 0.8" margins
    for section in doc.sections:
        section.page_width = Inches(8.5)
        section.page_height = Inches(11)
        section.left_margin = Inches(0.8)
        section.right_margin = Inches(0.8)
        section.top_margin = Inches(0.8)
        section.bottom_margin = Inches(0.8)

    # Theme colors
    theme = spec.get("theme", {})
    color_dark = theme.get("dark", "1A1A2E")
    color_accent = theme.get("accent", "E94560")
    color_accent2 = theme.get("accent2", "0F3460")
    color_light_bg = theme.get("lightBg", "F5F5F5")
    color_med_gray = theme.get("medGray", "666666")

    def _hex_rgb(hex_str):
        h = hex_str.lstrip("#")
        try:
            return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))
        except (ValueError, IndexError):
            return RGBColor(0, 0, 0)  # fallback to black on invalid hex

    def _set_cell_bg(cell, hex_color):
        """Set a table cell's background color."""
        h = hex_color.lstrip("#")
        shading = parse_xml(f'<w:shd {nsdecls("w")} w:fill="{h}"/>')
        cell._tc.get_or_add_tcPr().append(shading)

    # Header
    header_text = spec.get("header", "")
    if header_text:
        header = doc.sections[0].header
        hp = header.paragraphs[0]
        hp.text = header_text
        hp.style.font.size = Pt(9)
        hp.style.font.color.rgb = _hex_rgb(color_med_gray)
        hp.alignment = WD_ALIGN_PARAGRAPH.RIGHT

    # Footer
    footer_text = spec.get("footer", "")
    if footer_text:
        footer = doc.sections[0].footer
        fp = footer.paragraphs[0]
        fp.text = footer_text
        fp.style.font.size = Pt(9)
        fp.style.font.color.rgb = _hex_rgb(color_med_gray)
        fp.alignment = WD_ALIGN_PARAGRAPH.CENTER

    # Process sections or flat elements list
    # Supports two equivalent formats:
    #   {"elements": [...]}                           — flat list, no page breaks between sections
    #   {"sections": [{"children": [...]}, ...]}      — multi-section, page break between each
    flat_elements = spec.get("elements")
    if flat_elements is not None:
        sections_iter = [{"children": flat_elements}]
    else:
        sections_iter = spec.get("sections", [])

    for sec_idx, section_spec in enumerate(sections_iter):
        if sec_idx > 0:
            doc.add_page_break()

        for child in section_spec.get("children", section_spec.get("elements", [])):
            elem_type = child.get("type", "")

            if elem_type == "heading":
                level = child.get("level", 1)
                level = max(1, min(level, 6))
                p = doc.add_heading(child.get("text", ""), level=level)
                color = child.get("color", color_dark)
                for run in p.runs:
                    run.font.color.rgb = _hex_rgb(color)
                    run.font.name = "Arial"

            elif elem_type == "subtitle":
                p = doc.add_paragraph()
                run = p.add_run(child.get("text", ""))
                run.font.size = Pt(14)
                run.font.color.rgb = _hex_rgb(color_med_gray)
                run.font.name = "Arial"

            elif elem_type == "paragraph":
                p = doc.add_paragraph()
                runs = child.get("runs")
                if runs and isinstance(runs, list):
                    for run_spec in runs:
                        if isinstance(run_spec, str):
                            run_spec = {"text": run_spec}
                        run = p.add_run(run_spec.get("text", ""))
                        run.font.name = "Arial"
                        run.font.size = Pt(run_spec.get("fontSize", 11))
                        if run_spec.get("bold"):
                            run.bold = True
                        if run_spec.get("italic"):
                            run.italic = True
                        if run_spec.get("color"):
                            run.font.color.rgb = _hex_rgb(run_spec["color"])
                else:
                    text = child.get("text", "") or (runs if isinstance(runs, str) else "")
                    run = p.add_run(text)
                    run.font.name = "Arial"
                    if child.get("bold"):
                        run.bold = True
                    if child.get("italic"):
                        run.italic = True
                    if child.get("color"):
                        run.font.color.rgb = _hex_rgb(child["color"])

            elif elem_type == "accentBar":
                bar_color = child.get("color", color_accent)
                p = doc.add_paragraph()
                p.paragraph_format.space_before = Pt(0)
                p.paragraph_format.space_after = Pt(8)
                # Create a thin colored bar via a 1-row, 1-col table
                tbl = doc.add_table(rows=1, cols=1)
                tbl.alignment = WD_TABLE_ALIGNMENT.CENTER
                cell = tbl.cell(0, 0)
                cell.text = ""
                _set_cell_bg(cell, bar_color)
                # Set row height to ~4pt
                tr = tbl.rows[0]._tr
                trPr = tr.get_or_add_trPr()
                trHeight = parse_xml(f'<w:trHeight {nsdecls("w")} w:val="80" w:hRule="exact"/>')
                trPr.append(trHeight)

            elif elem_type == "divider":
                p = doc.add_paragraph()
                p.paragraph_format.space_before = Pt(6)
                p.paragraph_format.space_after = Pt(6)
                # Add bottom border
                pPr = p._p.get_or_add_pPr()
                pBdr = parse_xml(
                    f'<w:pBdr {nsdecls("w")}>'
                    f'  <w:bottom w:val="single" w:sz="4" w:space="1" w:color="DADCE0"/>'
                    f'</w:pBdr>'
                )
                pPr.append(pBdr)

            elif elem_type == "kpiRow":
                items = child.get("metrics") or child.get("items", [])
                if not items:
                    continue
                n = len(items)
                tbl = doc.add_table(rows=2, cols=n)
                tbl.alignment = WD_TABLE_ALIGNMENT.CENTER
                for i, kpi in enumerate(items):
                    bg = kpi.get("bg", color_accent2)
                    # Value cell
                    val_cell = tbl.cell(0, i)
                    val_cell.text = ""
                    vp = val_cell.paragraphs[0]
                    vr = vp.add_run(str(kpi.get("value", "")))
                    vr.font.size = Pt(22)
                    vr.font.bold = True
                    vr.font.color.rgb = RGBColor(255, 255, 255)
                    vr.font.name = "Arial"
                    vp.alignment = WD_ALIGN_PARAGRAPH.CENTER
                    _set_cell_bg(val_cell, bg)
                    # Label cell
                    lbl_cell = tbl.cell(1, i)
                    lbl_cell.text = ""
                    lp = lbl_cell.paragraphs[0]
                    lr = lp.add_run(str(kpi.get("label", "")))
                    lr.font.size = Pt(9)
                    lr.font.color.rgb = RGBColor(255, 255, 255)
                    lr.font.name = "Arial"
                    lp.alignment = WD_ALIGN_PARAGRAPH.CENTER
                    _set_cell_bg(lbl_cell, bg)
                # Remove table borders for clean KPI look
                for row in tbl.rows:
                    for cell in row.cells:
                        tc = cell._tc
                        tcPr = tc.get_or_add_tcPr()
                        tcBorders = parse_xml(
                            f'<w:tcBorders {nsdecls("w")}>'
                            f'  <w:top w:val="none" w:sz="0" w:space="0" w:color="auto"/>'
                            f'  <w:bottom w:val="none" w:sz="0" w:space="0" w:color="auto"/>'
                            f'  <w:left w:val="none" w:sz="0" w:space="0" w:color="auto"/>'
                            f'  <w:right w:val="none" w:sz="0" w:space="0" w:color="auto"/>'
                            f'</w:tcBorders>'
                        )
                        tcPr.append(tcBorders)

            elif elem_type == "callout":
                title_text = child.get("title", "")
                body_text = child.get("body", "")
                bg = child.get("bg", color_light_bg)
                tbl = doc.add_table(rows=1, cols=1)
                cell = tbl.cell(0, 0)
                cell.text = ""
                if title_text:
                    tp = cell.paragraphs[0]
                    tr = tp.add_run(title_text)
                    tr.font.bold = True
                    tr.font.size = Pt(12)
                    tr.font.name = "Arial"
                if body_text:
                    bp = cell.add_paragraph()
                    br = bp.add_run(body_text)
                    br.font.size = Pt(11)
                    br.font.name = "Arial"
                _set_cell_bg(cell, bg)

            elif elem_type == "table":
                headers = child.get("headers", [])
                rows = child.get("rows", [])
                style_spec = child.get("style", {})
                header_bg = style_spec.get("headerBg", color_accent2)
                zebra = style_spec.get("zebra", False)
                zebra_color = style_spec.get("zebraColor", "F5F5F5")

                all_rows = [headers] + rows if headers else rows
                if not all_rows:
                    continue

                n_cols = max(len(r) for r in all_rows)
                tbl = doc.add_table(rows=len(all_rows), cols=n_cols)
                tbl.alignment = WD_TABLE_ALIGNMENT.CENTER
                tbl.style = "Table Grid"

                for r_idx, row_data in enumerate(all_rows):
                    for c_idx, cell_val in enumerate(row_data):
                        if c_idx >= n_cols:
                            break
                        cell = tbl.cell(r_idx, c_idx)
                        cell.text = ""
                        p = cell.paragraphs[0]
                        run = p.add_run(str(cell_val))
                        run.font.name = "Arial"
                        run.font.size = Pt(10)

                        if r_idx == 0 and headers:
                            run.font.bold = True
                            run.font.color.rgb = RGBColor(255, 255, 255)
                            _set_cell_bg(cell, header_bg)
                        elif zebra and r_idx % 2 == 0:
                            _set_cell_bg(cell, zebra_color)

            elif elem_type == "richTable":
                headers = child.get("headers", [])
                rows = child.get("rows", [])
                style_spec = child.get("style", {})
                # Accept headerColor/zebraColor as top-level keys (flat spec format)
                # or headerBg/zebraColor inside nested "style" object (legacy format)
                header_bg = (child.get("headerColor")
                             or style_spec.get("headerBg", color_accent2))
                rich_zebra_color = (child.get("zebraColor")
                                    or style_spec.get("zebraColor"))

                all_rows = [headers] + rows if headers else rows
                if not all_rows:
                    continue

                n_cols = max(len(r) for r in all_rows)
                tbl = doc.add_table(rows=len(all_rows), cols=n_cols)
                tbl.alignment = WD_TABLE_ALIGNMENT.CENTER
                tbl.style = "Table Grid"

                for r_idx, row_data in enumerate(all_rows):
                    for c_idx, cell_val in enumerate(row_data):
                        if c_idx >= n_cols:
                            break
                        cell = tbl.cell(r_idx, c_idx)
                        cell.text = ""
                        p = cell.paragraphs[0]

                        if isinstance(cell_val, dict):
                            run = p.add_run(str(cell_val.get("text", "")))
                            run.font.name = "Arial"
                            run.font.size = Pt(cell_val.get("fontSize", 10))
                            if cell_val.get("bold"):
                                run.bold = True
                            if cell_val.get("italic"):
                                run.italic = True
                            if cell_val.get("color"):
                                run.font.color.rgb = _hex_rgb(cell_val["color"])
                            if cell_val.get("bg"):
                                _set_cell_bg(cell, cell_val["bg"])
                        else:
                            run = p.add_run(str(cell_val))
                            run.font.name = "Arial"
                            run.font.size = Pt(10)

                        if r_idx == 0 and headers:
                            run.font.bold = True
                            run.font.color.rgb = RGBColor(255, 255, 255)
                            _set_cell_bg(cell, header_bg)
                        elif rich_zebra_color and r_idx % 2 == 0:
                            _set_cell_bg(cell, rich_zebra_color)

            elif elem_type == "bullets":
                for item in child.get("items", []):
                    p = doc.add_paragraph(str(item), style="List Bullet")
                    for run in p.runs:
                        run.font.name = "Arial"
                        run.font.size = Pt(11)

            elif elem_type == "numberedList":
                for item in child.get("items", []):
                    p = doc.add_paragraph(str(item), style="List Number")
                    for run in p.runs:
                        run.font.name = "Arial"
                        run.font.size = Pt(11)

            elif elem_type == "personaCard":
                name = child.get("name", "")
                role = child.get("role", "")
                details = child.get("details", [])
                bg = child.get("bg", color_light_bg)

                tbl = doc.add_table(rows=1, cols=1)
                cell = tbl.cell(0, 0)
                cell.text = ""
                _set_cell_bg(cell, bg)

                # Name
                np = cell.paragraphs[0]
                nr = np.add_run(name)
                nr.font.bold = True
                nr.font.size = Pt(14)
                nr.font.name = "Arial"
                nr.font.color.rgb = _hex_rgb(color_dark)

                # Role
                if role:
                    rp = cell.add_paragraph()
                    rr = rp.add_run(role)
                    rr.font.size = Pt(11)
                    rr.font.italic = True
                    rr.font.name = "Arial"
                    rr.font.color.rgb = _hex_rgb(color_med_gray)

                # Details
                for detail in details:
                    dp = cell.add_paragraph()
                    dp.style = doc.styles["List Bullet"]
                    dr = dp.add_run(str(detail))
                    dr.font.size = Pt(10)
                    dr.font.name = "Arial"

            elif elem_type == "conceptHeader":
                title_text = child.get("title", "")
                sub_text = child.get("subtitle", "")
                color = child.get("color", color_dark)

                p = doc.add_paragraph()
                p.paragraph_format.space_before = Pt(18)
                p.paragraph_format.space_after = Pt(4)
                run = p.add_run(title_text.upper())
                run.font.bold = True
                run.font.size = Pt(13)
                run.font.color.rgb = _hex_rgb(color)
                run.font.name = "Arial"
                run.font.letter_spacing = Pt(1.5)

                if sub_text:
                    sp = doc.add_paragraph()
                    sp.paragraph_format.space_before = Pt(0)
                    sr = sp.add_run(sub_text)
                    sr.font.size = Pt(11)
                    sr.font.color.rgb = _hex_rgb(color_med_gray)
                    sr.font.name = "Arial"

            elif elem_type == "tag":
                text = child.get("text", "")
                tag_color = child.get("color", color_accent)
                tbl = doc.add_table(rows=1, cols=1)
                cell = tbl.cell(0, 0)
                cell.text = ""
                p = cell.paragraphs[0]
                run = p.add_run(f"  {text}  ")
                run.font.size = Pt(9)
                run.font.bold = True
                run.font.color.rgb = RGBColor(255, 255, 255)
                run.font.name = "Arial"
                _set_cell_bg(cell, tag_color)
                # Compact sizing
                tc = cell._tc
                tcPr = tc.get_or_add_tcPr()
                tcW = parse_xml(f'<w:tcW {nsdecls("w")} w:w="0" w:type="auto"/>')
                tcPr.append(tcW)

            elif elem_type == "pageBreak":
                doc.add_page_break()

            elif elem_type == "spacer":
                p = doc.add_paragraph()
                p.paragraph_format.space_before = Pt(child.get("height", 12))
                p.paragraph_format.space_after = Pt(0)

            elif elem_type == "image":
                img_path = child.get("path", "")
                width = child.get("width", 5.0)
                if img_path and os.path.exists(img_path):
                    doc.add_picture(img_path, width=Inches(width))

    doc.save(output_path)
    return output_path


def _upload_docx_as_doc(docx_path, title, folder_id=None):
    """Upload .docx to Drive, converting to Google Docs. Returns file metadata."""
    from googleapiclient.http import MediaFileUpload
    drive = _drive_service()

    metadata = {
        "name": title,
        "mimeType": "application/vnd.google-apps.document",
    }
    if folder_id:
        metadata["parents"] = [folder_id]

    media = MediaFileUpload(
        str(docx_path),
        mimetype="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        resumable=True,
    )
    uploaded = drive.files().create(
        body=metadata,
        media_body=media,
        fields="id, name, webViewLink",
    ).execute()
    return uploaded


def cmd_docs_build_rich(args):
    """Build a rich document from JSON spec, upload to Drive as Google Doc."""
    spec = _read_spec(args)
    title = args.title or spec.get("title", "Untitled Document")

    # Build .docx locally
    if args.output:
        output_path = pathlib.Path(args.output)
    else:
        output_path = pathlib.Path(tempfile.mkdtemp()) / f"{title}.docx"

    _build_rich_document(spec, str(output_path))

    if args.output:
        # Local-only mode — skip upload
        print(json.dumps({"localPath": str(output_path), "title": title}))
        return

    # Upload to Drive as Google Doc
    uploaded = _upload_docx_as_doc(str(output_path), title, folder_id=args.folder)
    doc_id = uploaded["id"]
    _auto_share_org(doc_id)
    url = uploaded.get("webViewLink", f"https://docs.google.com/document/d/{doc_id}/edit")

    result = {"documentId": doc_id, "url": url, "title": title, "orgShared": True}

    if args.link_share:
        result["linkShare"] = _link_share_file(doc_id)

    if args.share:
        drive = _drive_service()
        drive.permissions().create(
            fileId=doc_id,
            body={"type": "user", "role": "reader", "emailAddress": args.share},
            sendNotificationEmail=True,
            fields="id",
        ).execute()
        result["sharedWith"] = args.share

    # Clean up temp file
    output_path.unlink(missing_ok=True)

    print(json.dumps(result, indent=2))


def cmd_docs_build_rich_batch(args):
    """Build multiple rich documents in parallel from JSON spec array."""
    from concurrent.futures import ThreadPoolExecutor, as_completed

    specs = _read_spec(args)
    if not isinstance(specs, list):
        print("ERROR: Batch spec must be a JSON array of specs.", file=sys.stderr)
        sys.exit(1)

    max_workers = args.parallel or 4
    results = {"total": len(specs), "succeeded": 0, "failed": 0, "documents": []}

    def build_one(spec):
        title = spec.get("title", "Untitled Document")
        tmp_path = pathlib.Path(tempfile.mkdtemp()) / f"{title}.docx"
        try:
            _build_rich_document(spec, str(tmp_path))
            uploaded = _upload_docx_as_doc(str(tmp_path), title, folder_id=args.folder)
            doc_id = uploaded["id"]
            _auto_share_org(doc_id)
            url = uploaded.get("webViewLink", f"https://docs.google.com/document/d/{doc_id}/edit")
            entry = {"documentId": doc_id, "url": url, "title": title, "orgShared": True}
            if args.link_share:
                entry["linkShare"] = _link_share_file(doc_id)
            return entry
        except Exception as e:
            return {"title": title, "error": str(e)}
        finally:
            tmp_path.unlink(missing_ok=True)

    with ThreadPoolExecutor(max_workers=max_workers) as pool:
        futures = {pool.submit(build_one, s): s for s in specs}
        for future in as_completed(futures):
            entry = future.result()
            if "error" in entry:
                results["failed"] += 1
            else:
                results["succeeded"] += 1
            results["documents"].append(entry)

    print(json.dumps(results, indent=2))


def cmd_docs_replace_text(args):
    """Find and replace text in a document."""
    svc = _docs_service()
    request = {
        "replaceAllText": {
            "containsText": {
                "text": args.find,
                "matchCase": not args.ignore_case,
            },
            "replaceText": args.replace,
        }
    }
    resp = svc.documents().batchUpdate(
        documentId=args.document_id, body={"requests": [request]}
    ).execute()
    count = resp.get("replies", [{}])[0].get("replaceAllText", {}).get("occurrencesChanged", 0)
    print(json.dumps({"replaced": count, "find": args.find, "replaceWith": args.replace}))


def cmd_docs_delete_range(args):
    """Delete content between two indices in a document."""
    svc = _docs_service()

    start = args.start_index
    end = args.end_index

    if end <= start:
        print("ERROR: --end must be greater than --start", file=sys.stderr)
        sys.exit(1)

    svc.documents().batchUpdate(
        documentId=args.document_id,
        body={"requests": [{"deleteContentRange": {"range": {"startIndex": start, "endIndex": end}}}]},
    ).execute()
    print(json.dumps({"deleted": True, "startIndex": start, "endIndex": end, "chars": end - start}))


# ============================================================
# SHEETS commands
# ============================================================

def cmd_sheets_create(args):
    """Create a new spreadsheet."""
    svc = _sheets_service()
    body = {
        "properties": {"title": args.title},
        "sheets": [{"properties": {"title": "Sheet1"}}],
    }
    if args.frozen_rows:
        body["sheets"][0]["properties"]["gridProperties"] = {
            "frozenRowCount": args.frozen_rows
        }
    ss = svc.spreadsheets().create(body=body).execute()
    ss_id = ss["spreadsheetId"]
    _auto_share_org(ss_id)
    url = f"https://docs.google.com/spreadsheets/d/{ss_id}/edit"
    print(json.dumps({"spreadsheetId": ss_id, "url": url, "orgShared": True}))


def cmd_sheets_get(args):
    """Get spreadsheet metadata and sheet list."""
    svc = _sheets_service()
    ss = svc.spreadsheets().get(spreadsheetId=args.spreadsheet_id).execute()
    summary = {
        "spreadsheetId": ss["spreadsheetId"],
        "title": ss["properties"]["title"],
        "url": f"https://docs.google.com/spreadsheets/d/{ss['spreadsheetId']}/edit",
        "sheets": [
            {
                "sheetId": s["properties"]["sheetId"],
                "title": s["properties"]["title"],
                "rowCount": s["properties"].get("gridProperties", {}).get("rowCount"),
                "columnCount": s["properties"].get("gridProperties", {}).get("columnCount"),
            }
            for s in ss.get("sheets", [])
        ],
    }
    print(json.dumps(summary, indent=2))


def cmd_sheets_read(args):
    """Read data from a spreadsheet."""
    svc = _sheets_service()
    range_str = args.range or "Sheet1"
    render_map = {"formatted": "FORMATTED_VALUE", "raw": "UNFORMATTED_VALUE", "formula": "FORMULA"}
    render = render_map.get(args.render, "FORMATTED_VALUE")

    result = svc.spreadsheets().values().get(
        spreadsheetId=args.spreadsheet_id,
        range=range_str,
        valueRenderOption=render,
    ).execute()
    print(json.dumps(result.get("values", []), indent=2))


def cmd_sheets_write(args):
    """Write data to a spreadsheet."""
    svc = _sheets_service()
    data = _parse_json_arg(args.data, "data")
    input_option = "RAW" if args.input_type == "raw" else "USER_ENTERED"

    result = svc.spreadsheets().values().update(
        spreadsheetId=args.spreadsheet_id,
        range=args.range,
        valueInputOption=input_option,
        body={"values": data},
    ).execute()
    print(json.dumps({
        "updatedRange": result.get("updatedRange"),
        "updatedRows": result.get("updatedRows"),
        "updatedCells": result.get("updatedCells"),
    }))


def cmd_sheets_format(args):
    """Format cells in a spreadsheet."""
    svc = _sheets_service()
    style = _parse_json_arg(args.style, "style")

    # Parse range like A1:D1 to grid range
    grid_range = _parse_a1_to_grid(args.spreadsheet_id, args.range, svc)

    cell_format = {}
    if style.get("bold"):
        cell_format.setdefault("textFormat", {})["bold"] = True
    if style.get("italic"):
        cell_format.setdefault("textFormat", {})["italic"] = True
    if style.get("fontSize"):
        cell_format.setdefault("textFormat", {})["fontSize"] = style["fontSize"]
    if style.get("fontFamily"):
        cell_format.setdefault("textFormat", {})["fontFamily"] = style["fontFamily"]
    if style.get("textColor"):
        cell_format.setdefault("textFormat", {})["foregroundColorStyle"] = {
            "rgbColor": _hex_to_rgb(style["textColor"])
        }
    if style.get("backgroundColor"):
        cell_format["backgroundColorStyle"] = {
            "rgbColor": _hex_to_rgb(style["backgroundColor"])
        }
    if style.get("horizontalAlignment"):
        cell_format["horizontalAlignment"] = style["horizontalAlignment"]
    if style.get("numberFormat"):
        cell_format["numberFormat"] = style["numberFormat"]
    if style.get("wrapStrategy"):
        cell_format["wrapStrategy"] = style["wrapStrategy"]

    requests = [{
        "repeatCell": {
            "range": grid_range,
            "cell": {"userEnteredFormat": cell_format},
            "fields": "userEnteredFormat",
        }
    }]

    svc.spreadsheets().batchUpdate(
        spreadsheetId=args.spreadsheet_id, body={"requests": requests}
    ).execute()
    print(json.dumps({"formatted": args.range}))


def _parse_a1_to_grid(spreadsheet_id, a1_range, svc):
    """Parse A1 notation (e.g. 'Sheet1!A1:D5') to a GridRange dict."""
    import re

    sheet_title = "Sheet1"
    cell_range = a1_range
    if "!" in a1_range:
        sheet_title, cell_range = a1_range.split("!", 1)
        sheet_title = sheet_title.strip("'")

    # Get sheet ID from title
    ss = svc.spreadsheets().get(spreadsheetId=spreadsheet_id, fields="sheets.properties").execute()
    sheet_id = 0
    for s in ss.get("sheets", []):
        if s["properties"]["title"] == sheet_title:
            sheet_id = s["properties"]["sheetId"]
            break

    # Parse A1:B2 -> column/row indices
    match = re.match(r"([A-Z]+)(\d+)(?::([A-Z]+)(\d+))?", cell_range)
    if not match:
        return {"sheetId": sheet_id}

    def col_to_idx(col_str):
        idx = 0
        for c in col_str:
            idx = idx * 26 + (ord(c) - ord("A") + 1)
        return idx - 1

    start_col = col_to_idx(match.group(1))
    start_row = int(match.group(2)) - 1
    end_col = col_to_idx(match.group(3)) + 1 if match.group(3) else start_col + 1
    end_row = int(match.group(4)) if match.group(4) else start_row + 1

    return {
        "sheetId": sheet_id,
        "startRowIndex": start_row,
        "endRowIndex": end_row,
        "startColumnIndex": start_col,
        "endColumnIndex": end_col,
    }


def cmd_sheets_add_chart(args):
    """Add a native chart to a spreadsheet."""
    svc = _sheets_service()
    grid_range = _parse_a1_to_grid(args.spreadsheet_id, args.range, svc)

    chart_type_map = {
        "bar": "BAR", "line": "LINE", "pie": "PIE",
        "column": "COLUMN", "area": "AREA", "scatter": "SCATTER",
    }
    chart_type = chart_type_map.get(args.chart_type, "BAR")

    # Split range: first column = domain, rest = series
    start_col = grid_range["startColumnIndex"]
    end_col = grid_range["endColumnIndex"]

    domain_range = {**grid_range, "startColumnIndex": start_col, "endColumnIndex": start_col + 1}
    series_list = []
    for col in range(start_col + 1, end_col):
        series_list.append({
            "series": {"sourceRange": {"sources": [{**grid_range, "startColumnIndex": col, "endColumnIndex": col + 1}]}},
            "targetAxis": "BOTTOM_AXIS" if chart_type == "BAR" else "LEFT_AXIS",
        })

    chart_spec = {
        "basicChart": {
            "chartType": chart_type,
            "legendPosition": "BOTTOM_LEGEND",
            "domains": [{"domain": {"sourceRange": {"sources": [domain_range]}}}],
            "series": series_list,
            "headerCount": 1,
        }
    }

    if chart_type == "PIE":
        chart_spec = {
            "pieChart": {
                "legendPosition": "RIGHT_LEGEND",
                "domain": {"sourceRange": {"sources": [domain_range]}},
                "series": {"sourceRange": {"sources": [{
                    **grid_range,
                    "startColumnIndex": start_col + 1,
                    "endColumnIndex": start_col + 2,
                }]}},
            }
        }

    request = {
        "addChart": {
            "chart": {
                "spec": chart_spec,
                "position": {
                    "overlayPosition": {
                        "anchorCell": {
                            "sheetId": grid_range["sheetId"],
                            "rowIndex": grid_range.get("endRowIndex", 0) + 1,
                            "columnIndex": start_col,
                        },
                        "widthPixels": 600,
                        "heightPixels": 400,
                    }
                },
            }
        }
    }

    resp = svc.spreadsheets().batchUpdate(
        spreadsheetId=args.spreadsheet_id, body={"requests": [request]}
    ).execute()
    chart_id = resp["replies"][0]["addChart"]["chart"]["chartId"]
    print(json.dumps({"chartId": chart_id, "type": chart_type}))


def cmd_sheets_add_tab(args):
    """Add a new tab/sheet to a spreadsheet."""
    svc = _sheets_service()
    props = {"title": args.tab_title}
    if args.color:
        props["tabColorStyle"] = {"rgbColor": _hex_to_rgb(args.color)}

    resp = svc.spreadsheets().batchUpdate(
        spreadsheetId=args.spreadsheet_id,
        body={"requests": [{"addSheet": {"properties": props}}]},
    ).execute()
    sheet_id = resp["replies"][0]["addSheet"]["properties"]["sheetId"]
    print(json.dumps({"sheetId": sheet_id, "title": args.tab_title}))


def cmd_sheets_delete_tab(args):
    """Delete a tab/sheet from a spreadsheet."""
    svc = _sheets_service()

    sheet_id = args.sheet_id
    # If a title was given instead of numeric ID, resolve it
    if args.title:
        ss = svc.spreadsheets().get(spreadsheetId=args.spreadsheet_id).execute()
        found = None
        for s in ss.get("sheets", []):
            if s["properties"]["title"] == args.title:
                found = s["properties"]["sheetId"]
                break
        if found is None:
            print(f"ERROR: Tab '{args.title}' not found.", file=sys.stderr)
            sys.exit(1)
        sheet_id = found

    if sheet_id is None:
        print("ERROR: Provide --sheet-id or --title to identify the tab.", file=sys.stderr)
        sys.exit(1)

    svc.spreadsheets().batchUpdate(
        spreadsheetId=args.spreadsheet_id,
        body={"requests": [{"deleteSheet": {"sheetId": sheet_id}}]},
    ).execute()
    print(json.dumps({"deleted": True, "sheetId": sheet_id}))


def cmd_sheets_conditional_format(args):
    """Add conditional formatting rules."""
    svc = _sheets_service()
    rule = _parse_json_arg(args.rule, "rule")
    grid_range = _parse_a1_to_grid(args.spreadsheet_id, args.range, svc)

    request = {
        "addConditionalFormatRule": {
            "rule": {
                "ranges": [grid_range],
                **rule,
            },
            "index": 0,
        }
    }

    svc.spreadsheets().batchUpdate(
        spreadsheetId=args.spreadsheet_id, body={"requests": [request]}
    ).execute()
    print(json.dumps({"conditionalFormatAdded": True}))


def cmd_sheets_auto_resize(args):
    """Auto-resize columns to fit content."""
    svc = _sheets_service()

    # Parse column range like "0:10"
    start_col = 0
    end_col = 26
    if args.columns:
        parts = args.columns.split(":")
        start_col = int(parts[0])
        end_col = int(parts[1]) + 1 if len(parts) > 1 else start_col + 1

    # Get first sheet ID
    ss = svc.spreadsheets().get(spreadsheetId=args.spreadsheet_id, fields="sheets.properties").execute()
    sheet_id = ss["sheets"][0]["properties"]["sheetId"]

    svc.spreadsheets().batchUpdate(
        spreadsheetId=args.spreadsheet_id,
        body={"requests": [{
            "autoResizeDimensions": {
                "dimensions": {
                    "sheetId": sheet_id,
                    "dimension": "COLUMNS",
                    "startIndex": start_col,
                    "endIndex": end_col,
                }
            }
        }]},
    ).execute()
    print(json.dumps({"autoResized": True, "columns": f"{start_col}:{end_col - 1}"}))


# ============================================================
# DRIVE commands
# ============================================================

def cmd_drive_list(args):
    """List files in Drive."""
    drive = _drive_service()
    q_parts = ["trashed = false"]
    if args.folder:
        q_parts.append(f"'{args.folder}' in parents")
    if args.file_type:
        mime_map = {
            "slides": "application/vnd.google-apps.presentation",
            "docs": "application/vnd.google-apps.document",
            "sheets": "application/vnd.google-apps.spreadsheet",
        }
        if args.file_type in mime_map:
            q_parts.append(f"mimeType = '{mime_map[args.file_type]}'")

    results = drive.files().list(
        q=" and ".join(q_parts),
        pageSize=50,
        fields="files(id, name, mimeType, modifiedTime, webViewLink)",
        orderBy="modifiedTime desc",
    ).execute()

    files = results.get("files", [])
    print(json.dumps(files, indent=2))


def cmd_drive_share(args):
    """Share a file with a user."""
    drive = _drive_service()
    drive.permissions().create(
        fileId=args.file_id,
        body={"type": "user", "role": args.role, "emailAddress": args.email},
        sendNotificationEmail=True,
    ).execute()
    print(json.dumps({"shared": True, "email": args.email, "role": args.role}))


def cmd_drive_move(args):
    """Move a file to a different folder."""
    drive = _drive_service()
    # Get current parents
    f = drive.files().get(fileId=args.file_id, fields="parents").execute()
    prev_parents = ",".join(f.get("parents", []))

    drive.files().update(
        fileId=args.file_id,
        addParents=args.to,
        removeParents=prev_parents,
        fields="id, parents",
    ).execute()
    print(json.dumps({"moved": True, "newParent": args.to}))


def cmd_drive_upload(args):
    """Upload a local file to Drive."""
    from googleapiclient.http import MediaFileUpload
    drive = _drive_service()

    file_path = pathlib.Path(args.local_path)
    if not file_path.exists():
        print(f"ERROR: File not found: {args.local_path}", file=sys.stderr)
        sys.exit(1)

    metadata = {"name": file_path.name}
    if args.folder:
        metadata["parents"] = [args.folder]
    if args.mime_type:
        metadata["mimeType"] = args.mime_type

    media = MediaFileUpload(str(file_path), resumable=True)
    uploaded = drive.files().create(
        body=metadata,
        media_body=media,
        fields="id, name, webViewLink",
    ).execute()
    print(json.dumps(uploaded, indent=2))


def cmd_drive_create_folder(args):
    """Create a folder in Drive."""
    drive = _drive_service()
    metadata = {
        "name": args.name,
        "mimeType": "application/vnd.google-apps.folder",
    }
    if args.parent:
        metadata["parents"] = [args.parent]

    folder = drive.files().create(body=metadata, fields="id, name, webViewLink").execute()
    print(json.dumps(folder, indent=2))


def cmd_drive_rename(args):
    """Rename a file or folder in Drive."""
    drive = _drive_service()
    updated = drive.files().update(
        fileId=args.file_id,
        body={"name": args.new_name},
        fields="id, name, webViewLink",
    ).execute()
    print(json.dumps(updated, indent=2))


def cmd_drive_search(args):
    """Search files in Drive by name query."""
    drive = _drive_service()
    q_parts = ["trashed = false"]
    q_parts.append(f"name contains '{args.query}'")
    if args.parent:
        q_parts.append(f"'{args.parent}' in parents")
    if args.folder_only:
        q_parts.append("mimeType = 'application/vnd.google-apps.folder'")

    limit = args.limit or 25
    results = drive.files().list(
        q=" and ".join(q_parts),
        pageSize=limit,
        fields="files(id, name, mimeType, modifiedTime, webViewLink, parents)",
        orderBy="modifiedTime desc",
    ).execute()

    files = results.get("files", [])
    print(json.dumps(files, indent=2))


# ============================================================
# Rich XLSX pipeline (openpyxl + pandas)
# ============================================================

def _link_share_file(file_id):
    """Make a Drive file link-shareable (anyone with link can view). Returns share URL."""
    drive = _drive_service()
    drive.permissions().create(
        fileId=file_id,
        body={"type": "anyone", "role": "reader"},
        fields="id",
    ).execute()
    info = drive.files().get(fileId=file_id, fields="webViewLink").execute()
    return info["webViewLink"]


def _col_letter_to_idx(col_str):
    """Convert column letter(s) to 0-based index. A=0, B=1, Z=25, AA=26."""
    idx = 0
    for c in col_str.upper():
        idx = idx * 26 + (ord(c) - ord("A") + 1)
    return idx - 1


def _parse_range_to_cells(range_str):
    """Parse 'A1:D5' into ((row_start, col_start), (row_end, col_end)) as 1-based."""
    import re
    m = re.match(r"([A-Z]+)(\d+)(?::([A-Z]+)(\d+))?", range_str.upper())
    if not m:
        return None
    c1 = _col_letter_to_idx(m.group(1)) + 1
    r1 = int(m.group(2))
    if m.group(3):
        c2 = _col_letter_to_idx(m.group(3)) + 1
        r2 = int(m.group(4))
    else:
        c2, r2 = c1, r1
    return ((r1, c1), (r2, c2))


def _hex_to_openpyxl_color(hex_str):
    """Convert hex color (with or without #) to openpyxl aRGB string."""
    h = hex_str.lstrip("#")
    if len(h) == 6:
        return "FF" + h.upper()
    if len(h) == 8:
        return h.upper()
    return "FF000000"


def _build_rich_spreadsheet(spec, output_path):
    """Build a rich .xlsx file from a JSON spec using openpyxl.

    spec format:
    {
      "title": "...",
      "sheets": [{
        "name": "Sheet1",
        "tabColor": "1a73e8",
        "frozenRows": 1, "frozenCols": 0,
        "autoFilter": "A1:D1",
        "columnWidths": {"A": 25, "B": 15},
        "data": [[...], [...]],
        "formulas": {"D2": "=B2*C2"},
        "styles": [{"range": "A1:D1", "bold": true, ...}],
        "conditionalFormats": [{"range": "...", "type": "greaterThan", ...}],
        "charts": [{"type": "bar", "title": "...", "dataRange": "A1:B4", ...}]
      }]
    }
    """
    from openpyxl import Workbook
    from openpyxl.styles import Font, PatternFill, Alignment, Border, Side, numbers
    from openpyxl.chart import BarChart, LineChart, PieChart, AreaChart, ScatterChart, Reference
    from openpyxl.formatting.rule import CellIsRule, ColorScaleRule, DataBarRule, FormulaRule
    from openpyxl.utils import get_column_letter

    wb = Workbook()

    # Professional defaults
    default_font = Font(name="Arial", size=10)
    thin_border_side = Side(style="thin", color="FFDADCE0")

    sheets_spec = spec.get("sheets", [])
    if not sheets_spec:
        sheets_spec = [{"name": "Sheet1", "data": []}]

    for i, sheet_spec in enumerate(sheets_spec):
        if i == 0:
            ws = wb.active
            ws.title = sheet_spec.get("name", "Sheet1")
        else:
            ws = wb.create_sheet(title=sheet_spec.get("name", f"Sheet{i+1}"))

        # Tab color
        if sheet_spec.get("tabColor"):
            ws.sheet_properties.tabColor = sheet_spec["tabColor"].lstrip("#")

        # Write data
        data = sheet_spec.get("data", [])
        for r_idx, row in enumerate(data, 1):
            for c_idx, val in enumerate(row, 1):
                cell = ws.cell(row=r_idx, column=c_idx)
                if val is not None:
                    cell.value = val
                cell.font = default_font

        # Apply formulas
        for cell_ref, formula in sheet_spec.get("formulas", {}).items():
            ws[cell_ref] = formula

        # Apply thin borders on data range
        if data:
            max_row = len(data)
            max_col = max(len(row) for row in data) if data else 0
            thin_border = Border(
                left=thin_border_side, right=thin_border_side,
                top=thin_border_side, bottom=thin_border_side,
            )
            for r in range(1, max_row + 1):
                for c in range(1, max_col + 1):
                    ws.cell(row=r, column=c).border = thin_border

        # Apply styles
        for style_spec in sheet_spec.get("styles", []):
            rng = style_spec.get("range")
            if not rng:
                continue
            parsed = _parse_range_to_cells(rng)
            if not parsed:
                continue
            (r1, c1), (r2, c2) = parsed

            # Build style objects
            font_kwargs = {"name": "Arial", "size": 10}
            if style_spec.get("bold"):
                font_kwargs["bold"] = True
            if style_spec.get("italic"):
                font_kwargs["italic"] = True
            if style_spec.get("fontSize"):
                font_kwargs["size"] = style_spec["fontSize"]
            if style_spec.get("fontFamily"):
                font_kwargs["name"] = style_spec["fontFamily"]
            if style_spec.get("textColor"):
                font_kwargs["color"] = _hex_to_openpyxl_color(style_spec["textColor"])

            fill = None
            if style_spec.get("backgroundColor"):
                fill = PatternFill(
                    start_color=_hex_to_openpyxl_color(style_spec["backgroundColor"]),
                    end_color=_hex_to_openpyxl_color(style_spec["backgroundColor"]),
                    fill_type="solid",
                )

            alignment = None
            align_kwargs = {}
            if style_spec.get("alignment"):
                align_kwargs["horizontal"] = style_spec["alignment"]
            if style_spec.get("wrapText"):
                align_kwargs["wrap_text"] = True
            if align_kwargs:
                alignment = Alignment(**align_kwargs)

            num_fmt = style_spec.get("numberFormat")

            font = Font(**font_kwargs)
            for r in range(r1, r2 + 1):
                for c in range(c1, c2 + 1):
                    cell = ws.cell(row=r, column=c)
                    cell.font = font
                    if fill:
                        cell.fill = fill
                    if alignment:
                        cell.alignment = alignment
                    if num_fmt:
                        cell.number_format = num_fmt

        # Column widths
        for col_letter, width in sheet_spec.get("columnWidths", {}).items():
            ws.column_dimensions[col_letter.upper()].width = width

        # Frozen rows/cols
        frozen_rows = sheet_spec.get("frozenRows", 0)
        frozen_cols = sheet_spec.get("frozenCols", 0)
        if frozen_rows or frozen_cols:
            freeze_cell = f"{get_column_letter(frozen_cols + 1)}{frozen_rows + 1}"
            ws.freeze_panes = freeze_cell

        # Auto-filter
        if sheet_spec.get("autoFilter"):
            ws.auto_filter.ref = sheet_spec["autoFilter"]

        # Conditional formats
        for cf_spec in sheet_spec.get("conditionalFormats", []):
            cf_range = cf_spec.get("range")
            cf_type = cf_spec.get("type")
            if not cf_range or not cf_type:
                continue

            cf_fill = None
            cf_font = None
            if cf_spec.get("fill"):
                cf_fill = PatternFill(
                    start_color=_hex_to_openpyxl_color(cf_spec["fill"]),
                    end_color=_hex_to_openpyxl_color(cf_spec["fill"]),
                    fill_type="solid",
                )
            if cf_spec.get("fontColor"):
                cf_font = Font(color=_hex_to_openpyxl_color(cf_spec["fontColor"]))

            if cf_type == "colorScale":
                rule = ColorScaleRule(
                    start_type="min", start_color=_hex_to_openpyxl_color(cf_spec.get("startColor", "FF0000")),
                    end_type="max", end_color=_hex_to_openpyxl_color(cf_spec.get("endColor", "00FF00")),
                )
                ws.conditional_formatting.add(cf_range, rule)
            elif cf_type == "dataBar":
                rule = DataBarRule(
                    start_type="min", end_type="max",
                    color=_hex_to_openpyxl_color(cf_spec.get("color", "638EC6")),
                )
                ws.conditional_formatting.add(cf_range, rule)
            elif cf_type in ("greaterThan", "lessThan", "equal", "between"):
                op_map = {
                    "greaterThan": "greaterThan",
                    "lessThan": "lessThan",
                    "equal": "equal",
                    "between": "between",
                }
                val = cf_spec.get("value")
                if cf_type == "between":
                    formula = [cf_spec.get("start", 0), cf_spec.get("end", 0)]
                else:
                    formula = [val]
                rule = CellIsRule(
                    operator=op_map[cf_type],
                    formula=formula,
                    fill=cf_fill,
                    font=cf_font,
                )
                ws.conditional_formatting.add(cf_range, rule)
            elif cf_type == "containsText":
                text = cf_spec.get("text", "")
                rule = FormulaRule(
                    formula=[f'NOT(ISERROR(SEARCH("{text}",{cf_range.split(":")[0]})))'],
                    fill=cf_fill,
                    font=cf_font,
                )
                ws.conditional_formatting.add(cf_range, rule)

        # Charts
        for chart_spec in sheet_spec.get("charts", []):
            chart_type = chart_spec.get("type", "bar")
            chart_title = chart_spec.get("title", "")
            data_range = chart_spec.get("dataRange", "")
            position = chart_spec.get("position", "F2")
            width = chart_spec.get("width", 15)
            height = chart_spec.get("height", 10)

            if not data_range:
                continue
            parsed = _parse_range_to_cells(data_range)
            if not parsed:
                continue
            (dr1, dc1), (dr2, dc2) = parsed

            chart_class_map = {
                "bar": BarChart, "line": LineChart, "pie": PieChart,
                "area": AreaChart, "scatter": ScatterChart,
            }
            ChartClass = chart_class_map.get(chart_type, BarChart)
            chart_obj = ChartClass()
            chart_obj.title = chart_title
            chart_obj.width = width
            chart_obj.height = height

            if chart_type == "pie":
                data_ref = Reference(ws, min_col=dc1 + 1, min_row=dr1, max_row=dr2)
                cats = Reference(ws, min_col=dc1, min_row=dr1 + 1, max_row=dr2)
                chart_obj.add_data(data_ref, titles_from_data=True)
                chart_obj.set_categories(cats)
            else:
                data_ref = Reference(ws, min_col=dc1 + 1, max_col=dc2, min_row=dr1, max_row=dr2)
                cats = Reference(ws, min_col=dc1, min_row=dr1 + 1, max_row=dr2)
                chart_obj.add_data(data_ref, titles_from_data=True)
                chart_obj.set_categories(cats)

            ws.add_chart(chart_obj, position)

    wb.save(output_path)
    return output_path


def _upload_xlsx_as_sheets(xlsx_path, title, folder_id=None):
    """Upload .xlsx to Drive, converting to Google Sheets. Returns file metadata."""
    from googleapiclient.http import MediaFileUpload
    drive = _drive_service()

    metadata = {
        "name": title,
        "mimeType": "application/vnd.google-apps.spreadsheet",
    }
    if folder_id:
        metadata["parents"] = [folder_id]

    media = MediaFileUpload(str(xlsx_path), mimetype="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet", resumable=True)
    uploaded = drive.files().create(
        body=metadata,
        media_body=media,
        fields="id, name, webViewLink",
    ).execute()
    return uploaded


def _read_spec(args):
    """Read JSON spec from --spec (inline JSON), --spec-file, or stdin."""
    if hasattr(args, "spec") and args.spec:
        try:
            return json.loads(args.spec)
        except json.JSONDecodeError as e:
            print(f"ERROR: Invalid JSON for --spec: {e}", file=sys.stderr)
            sys.exit(1)
    if hasattr(args, "spec_file") and args.spec_file:
        p = pathlib.Path(args.spec_file)
        if not p.exists():
            print(f"ERROR: Spec file not found: {args.spec_file}", file=sys.stderr)
            sys.exit(1)
        return json.loads(p.read_text())
    else:
        data = sys.stdin.read()
        if not data.strip():
            print("ERROR: No spec provided. Use --spec, --spec-file, or pipe JSON to stdin.", file=sys.stderr)
            sys.exit(1)
        try:
            return json.loads(data)
        except json.JSONDecodeError as e:
            print(f"ERROR: Invalid JSON from stdin: {e}", file=sys.stderr)
            sys.exit(1)


def cmd_sheets_build_rich(args):
    """Build a rich spreadsheet from JSON spec, upload to Drive as Google Sheets."""
    spec = _read_spec(args)
    title = args.title or spec.get("title", "Untitled Spreadsheet")

    # Build .xlsx locally
    if args.output:
        output_path = pathlib.Path(args.output)
    else:
        output_path = pathlib.Path(tempfile.mkdtemp()) / f"{title}.xlsx"

    _build_rich_spreadsheet(spec, str(output_path))

    # Upload to Drive as Google Sheets
    uploaded = _upload_xlsx_as_sheets(str(output_path), title, folder_id=args.folder)
    ss_id = uploaded["id"]
    _auto_share_org(ss_id)
    url = uploaded.get("webViewLink", f"https://docs.google.com/spreadsheets/d/{ss_id}/edit")

    result = {"spreadsheetId": ss_id, "url": url, "title": title, "orgShared": True}

    if args.link_share:
        share_url = _link_share_file(ss_id)
        result["linkShare"] = share_url

    if hasattr(args, "share") and args.share:
        drive = _drive_service()
        drive.permissions().create(
            fileId=ss_id,
            body={"type": "user", "role": "reader", "emailAddress": args.share},
            sendNotificationEmail=True,
            fields="id",
        ).execute()
        result["sharedWith"] = args.share

    # Clean up temp file if we created one
    if not args.output:
        output_path.unlink(missing_ok=True)

    print(json.dumps(result, indent=2))


def cmd_sheets_build_rich_batch(args):
    """Build multiple rich spreadsheets in parallel from JSON spec array."""
    from concurrent.futures import ThreadPoolExecutor, as_completed

    specs = _read_spec(args)
    if not isinstance(specs, list):
        print("ERROR: Batch spec must be a JSON array of specs.", file=sys.stderr)
        sys.exit(1)

    max_workers = args.parallel or 4
    results = {"total": len(specs), "succeeded": 0, "failed": 0, "spreadsheets": []}

    def build_one(spec):
        title = spec.get("title", "Untitled Spreadsheet")
        tmp_path = pathlib.Path(tempfile.mkdtemp()) / f"{title}.xlsx"
        try:
            _build_rich_spreadsheet(spec, str(tmp_path))
            uploaded = _upload_xlsx_as_sheets(str(tmp_path), title, folder_id=args.folder)
            ss_id = uploaded["id"]
            _auto_share_org(ss_id)
            url = uploaded.get("webViewLink", f"https://docs.google.com/spreadsheets/d/{ss_id}/edit")
            entry = {"spreadsheetId": ss_id, "url": url, "title": title, "orgShared": True}
            if args.link_share:
                entry["linkShare"] = _link_share_file(ss_id)
            return entry
        except Exception as e:
            return {"title": title, "error": str(e)}
        finally:
            tmp_path.unlink(missing_ok=True)

    with ThreadPoolExecutor(max_workers=max_workers) as pool:
        futures = {pool.submit(build_one, s): s for s in specs}
        for future in as_completed(futures):
            entry = future.result()
            if "error" in entry:
                results["failed"] += 1
            else:
                results["succeeded"] += 1
            results["spreadsheets"].append(entry)

    print(json.dumps(results, indent=2))


def cmd_sheets_analyze(args):
    """Analyze a spreadsheet's structure and data statistics."""
    import pandas as pd
    from openpyxl import load_workbook

    # Get the .xlsx file
    if args.file:
        xlsx_path = args.file
    else:
        # Download from Google Sheets
        xlsx_path = os.path.join(tempfile.mkdtemp(), "analyze.xlsx")
        _download_sheet_as_xlsx(args.spreadsheet_id, xlsx_path)

    # openpyxl analysis: structure
    wb = load_workbook(xlsx_path, data_only=True)
    structure = {
        "sheetNames": wb.sheetnames,
        "sheets": [],
    }

    target_sheets = [args.sheet] if args.sheet else wb.sheetnames
    for sheet_name in target_sheets:
        if sheet_name not in wb.sheetnames:
            continue
        ws = wb[sheet_name]
        sheet_info = {
            "name": sheet_name,
            "dimensions": ws.dimensions,
            "maxRow": ws.max_row,
            "maxColumn": ws.max_column,
            "mergedCells": [str(m) for m in ws.merged_cells.ranges],
            "frozenPanes": str(ws.freeze_panes) if ws.freeze_panes else None,
        }

        # Find formulas (reload without data_only)
        wb_formulas = load_workbook(xlsx_path, data_only=False)
        ws_f = wb_formulas[sheet_name]
        formulas = {}
        for row in ws_f.iter_rows():
            for cell in row:
                if cell.value and isinstance(cell.value, str) and cell.value.startswith("="):
                    formulas[cell.coordinate] = cell.value
        sheet_info["formulaCount"] = len(formulas)
        if formulas:
            # Show first 20
            sheet_info["formulaSample"] = dict(list(formulas.items())[:20])

        structure["sheets"].append(sheet_info)
    wb.close()

    # pandas analysis: data stats
    pandas_info = {}
    for sheet_name in target_sheets:
        try:
            df = pd.read_excel(xlsx_path, sheet_name=sheet_name)
            pandas_info[sheet_name] = {
                "shape": list(df.shape),
                "columns": list(df.columns),
                "dtypes": {col: str(dtype) for col, dtype in df.dtypes.items()},
                "nullCounts": {col: int(cnt) for col, cnt in df.isnull().sum().items() if cnt > 0},
                "head": df.head(5).to_dict(orient="records"),
                "describe": json.loads(df.describe(include="all").to_json()),
            }
        except Exception as e:
            pandas_info[sheet_name] = {"error": str(e)}

    result = {"structure": structure, "data": pandas_info}
    if not args.file:
        os.unlink(xlsx_path)

    print(json.dumps(result, indent=2, default=str))


def _download_sheet_as_xlsx(spreadsheet_id, output_path):
    """Download a Google Sheet as .xlsx via Drive export."""
    drive = _drive_service()
    request = drive.files().export_media(
        fileId=spreadsheet_id,
        mimeType="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    )
    with open(output_path, "wb") as f:
        from googleapiclient.http import MediaIoBaseDownload
        downloader = MediaIoBaseDownload(f, request)
        done = False
        while not done:
            _, done = downloader.next_chunk()


def cmd_sheets_query(args):
    """Query/filter/aggregate data from a spreadsheet using pandas."""
    import pandas as pd

    # Get the .xlsx file
    if args.file:
        xlsx_path = args.file
    else:
        xlsx_path = os.path.join(tempfile.mkdtemp(), "query.xlsx")
        _download_sheet_as_xlsx(args.spreadsheet_id, xlsx_path)

    sheet_name = args.sheet or 0
    df = pd.read_excel(xlsx_path, sheet_name=sheet_name)

    # Apply filter
    if args.filter:
        filters = _parse_json_arg(args.filter, "filter")
        for col, condition in filters.items():
            if isinstance(condition, dict):
                for op, val in condition.items():
                    if op == ">":
                        df = df[df[col] > val]
                    elif op == "<":
                        df = df[df[col] < val]
                    elif op == ">=":
                        df = df[df[col] >= val]
                    elif op == "<=":
                        df = df[df[col] <= val]
                    elif op == "!=":
                        df = df[df[col] != val]
                    elif op == "contains":
                        df = df[df[col].astype(str).str.contains(str(val), case=False, na=False)]
            else:
                df = df[df[col] == condition]

    # Apply groupby + agg
    if args.groupby:
        agg_spec = _parse_json_arg(args.agg, "agg") if args.agg else "count"
        df = df.groupby(args.groupby).agg(agg_spec).reset_index()

    # Sort
    if args.sort:
        ascending = not args.sort.startswith("-")
        sort_col = args.sort.lstrip("-")
        df = df.sort_values(sort_col, ascending=ascending)

    # Limit
    if args.limit:
        df = df.head(args.limit)

    result = {
        "rowCount": len(df),
        "columns": list(df.columns),
        "data": json.loads(df.to_json(orient="records", date_format="iso", default_handler=str)),
    }

    if not args.file:
        os.unlink(xlsx_path)

    print(json.dumps(result, indent=2))


def cmd_sheets_download(args):
    """Download a Google Sheet as .xlsx."""
    output = args.output or f"/tmp/{args.spreadsheet_id}.xlsx"
    _download_sheet_as_xlsx(args.spreadsheet_id, output)
    print(json.dumps({"downloaded": output, "spreadsheetId": args.spreadsheet_id}))


def cmd_sheets_edit(args):
    """Edit an existing Google Sheet: download -> apply edits -> re-upload."""
    from openpyxl import load_workbook
    from openpyxl.styles import Font, PatternFill, Alignment, Border, Side

    spec = _read_spec(args)
    ss_id = args.spreadsheet_id

    # Download current sheet
    tmp_dir = tempfile.mkdtemp()
    xlsx_path = os.path.join(tmp_dir, "edit.xlsx")
    _download_sheet_as_xlsx(ss_id, xlsx_path)

    # Load workbook
    wb = load_workbook(xlsx_path)

    # Apply edits per sheet (default to first sheet if no sheet name given)
    target_sheets = spec.get("sheets", [spec]) if "sheets" in spec else [spec]
    for edit_spec in target_sheets:
        sheet_name = edit_spec.get("name")
        if sheet_name and sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
        else:
            ws = wb.active

        # Write new data if provided
        for cell_ref, value in edit_spec.get("data", {}).items():
            ws[cell_ref] = value

        # Apply formulas
        for cell_ref, formula in edit_spec.get("formulas", {}).items():
            ws[cell_ref] = formula

        # Apply styles
        for style_spec in edit_spec.get("styles", []):
            rng = style_spec.get("range")
            if not rng:
                continue
            parsed = _parse_range_to_cells(rng)
            if not parsed:
                continue
            (r1, c1), (r2, c2) = parsed

            font_kwargs = {}
            if style_spec.get("bold"):
                font_kwargs["bold"] = True
            if style_spec.get("italic"):
                font_kwargs["italic"] = True
            if style_spec.get("fontSize"):
                font_kwargs["size"] = style_spec["fontSize"]
            if style_spec.get("fontFamily"):
                font_kwargs["name"] = style_spec["fontFamily"]
            if style_spec.get("textColor"):
                font_kwargs["color"] = _hex_to_openpyxl_color(style_spec["textColor"])

            fill = None
            if style_spec.get("backgroundColor"):
                fill = PatternFill(
                    start_color=_hex_to_openpyxl_color(style_spec["backgroundColor"]),
                    end_color=_hex_to_openpyxl_color(style_spec["backgroundColor"]),
                    fill_type="solid",
                )

            alignment = None
            align_kwargs = {}
            if style_spec.get("alignment"):
                align_kwargs["horizontal"] = style_spec["alignment"]
            if style_spec.get("wrapText"):
                align_kwargs["wrap_text"] = True
            if align_kwargs:
                alignment = Alignment(**align_kwargs)

            num_fmt = style_spec.get("numberFormat")

            for r in range(r1, r2 + 1):
                for c in range(c1, c2 + 1):
                    cell = ws.cell(row=r, column=c)
                    if font_kwargs:
                        cell.font = Font(**font_kwargs)
                    if fill:
                        cell.fill = fill
                    if alignment:
                        cell.alignment = alignment
                    if num_fmt:
                        cell.number_format = num_fmt

    # Save edited workbook
    edited_path = os.path.join(tmp_dir, "edited.xlsx")
    wb.save(edited_path)
    wb.close()

    # Delete old file on Drive, upload new one (preserving ID is not possible with Drive API
    # when converting, so we upload as new and keep the same title)
    drive = _drive_service()
    old_info = drive.files().get(fileId=ss_id, fields="name, parents").execute()
    title = old_info.get("name", "Edited Spreadsheet")
    parents = old_info.get("parents", [])

    # Trash the old file
    drive.files().update(fileId=ss_id, body={"trashed": True}).execute()

    # Upload edited version
    uploaded = _upload_xlsx_as_sheets(edited_path, title, folder_id=parents[0] if parents else None)
    new_id = uploaded["id"]
    _auto_share_org(new_id)
    url = uploaded.get("webViewLink", f"https://docs.google.com/spreadsheets/d/{new_id}/edit")

    result = {"spreadsheetId": new_id, "url": url, "title": title, "previousId": ss_id, "orgShared": True}
    if args.link_share:
        result["linkShare"] = _link_share_file(new_id)

    # Cleanup
    os.unlink(xlsx_path)
    os.unlink(edited_path)

    print(json.dumps(result, indent=2))


# ============================================================
# Argument parser
# ============================================================

def build_parser():
    parser = argparse.ArgumentParser(
        description="Google Workspace API helper (Slides, Docs, Sheets, Drive)",
        formatter_class=argparse.RawDescriptionHelpFormatter,
    )
    sub = parser.add_subparsers(dest="service", help="Service to interact with")

    # --- AUTH ---
    auth_p = sub.add_parser("auth", help="Authentication management")
    auth_sub = auth_p.add_subparsers(dest="auth_cmd")
    auth_sub.add_parser("status", help="Check credential status")

    # --- SLIDES ---
    slides_p = sub.add_parser("slides", help="Google Slides operations")
    slides_sub = slides_p.add_subparsers(dest="slides_cmd")

    # slides create
    p = slides_sub.add_parser("create", help="Create a presentation")
    p.add_argument("title", help="Presentation title")
    p.add_argument("--clean", action="store_true",
                   help="Remove default placeholder slide (recommended for layout-based decks)")

    # slides get
    p = slides_sub.add_parser("get", help="Get presentation structure")
    p.add_argument("presentation_id", help="Presentation ID")

    # slides add-slide
    p = slides_sub.add_parser("add-slide", help="Add a slide")
    p.add_argument("presentation_id", help="Presentation ID")
    p.add_argument("--layout", default="BLANK",
                   choices=["BLANK", "TITLE", "TITLE_AND_BODY", "SECTION_HEADER",
                            "TITLE_ONLY", "ONE_COLUMN_TEXT", "MAIN_POINT", "BIG_NUMBER",
                            "CAPTION_ONLY", "TITLE_AND_TWO_COLUMNS"])
    p.add_argument("--index", type=int, default=None, help="Insertion index")

    # slides add-text
    p = slides_sub.add_parser("add-text", help="Add a text box")
    p.add_argument("presentation_id")
    p.add_argument("slide_index", type=int)
    p.add_argument("text")
    p.add_argument("--style", help="JSON style object")
    p.add_argument("--heading", choices=["heading1", "heading2", "heading3"])

    # slides add-table
    p = slides_sub.add_parser("add-table", help="Add a styled table")
    p.add_argument("presentation_id")
    p.add_argument("slide_index", type=int)
    p.add_argument("--data", required=True, help="2D JSON array")
    p.add_argument("--style", help="JSON style object")

    # slides add-image
    p = slides_sub.add_parser("add-image", help="Add an image from URL")
    p.add_argument("presentation_id")
    p.add_argument("slide_index", type=int)
    p.add_argument("--url", required=True)
    p.add_argument("--x", type=int, help="X position in EMU")
    p.add_argument("--y", type=int, help="Y position in EMU")
    p.add_argument("--width", type=int, help="Width in EMU")
    p.add_argument("--height", type=int, help="Height in EMU")

    # slides add-shape
    p = slides_sub.add_parser("add-shape", help="Add a shape")
    p.add_argument("presentation_id")
    p.add_argument("slide_index", type=int)
    p.add_argument("--type", dest="shape_type", default="RECTANGLE",
                   help="Shape type (TEXT_BOX, RECTANGLE, ELLIPSE, etc.)")
    p.add_argument("--style", help="JSON style object")

    # slides set-background
    p = slides_sub.add_parser("set-background", help="Set slide background")
    p.add_argument("presentation_id")
    p.add_argument("slide_index", type=int)
    p.add_argument("--color", help="Hex color (#RRGGBB)")
    p.add_argument("--image", help="Image URL")

    # slides apply-theme
    p = slides_sub.add_parser("apply-theme", help="Apply brand theme")
    p.add_argument("presentation_id")
    p.add_argument("--set-backgrounds", action="store_true",
                   help="Also set slide backgrounds (skip if using add-layout)")

    # slides add-chart
    p = slides_sub.add_parser("add-chart", help="Add a chart")
    p.add_argument("presentation_id")
    p.add_argument("slide_index", type=int)
    p.add_argument("--data", required=True, help="JSON chart data")
    p.add_argument("--type", dest="chart_type", default="bar",
                   choices=["bar", "line", "pie", "column", "area", "scatter"])
    p.add_argument("--chart-source", dest="chart_source", default="image",
                   choices=["image", "sheets"])

    # slides add-layout
    p = slides_sub.add_parser("add-layout", help="Add a pre-designed professional slide layout")
    p.add_argument("presentation_id")
    p.add_argument("layout_type", choices=["title", "section", "kpi", "two-column", "closing", "content", "ad-card"])
    p.add_argument("--data", help="JSON data for layout (title, subtitle, kpis, body)")

    # slides export
    p = slides_sub.add_parser("export", help="Export presentation")
    p.add_argument("presentation_id")
    p.add_argument("--format", required=True, choices=["pdf", "pptx", "png", "svg"])
    p.add_argument("--output", help="Output file path")

    # slides batch-update
    p = slides_sub.add_parser("batch-update", help="Raw batchUpdate")
    p.add_argument("presentation_id")
    p.add_argument("--requests", required=True, help="JSON array of requests")

    # slides build-rich
    p = slides_sub.add_parser("build-rich", help="Build rich presentation from JSON spec")
    p.add_argument("--title", help="Presentation title (overrides spec)")
    p.add_argument("--folder", help="Drive folder ID")
    p.add_argument("--link-share", action="store_true", help="Make link-shareable")
    p.add_argument("--share", help="Email to share with")
    p.add_argument("--output", help="Save .pptx locally instead of uploading")
    p.add_argument("--spec", help="Inline JSON spec")
    p.add_argument("--spec-file", help="Path to JSON spec file")
    p.add_argument("--qa", action="store_true", help="Run vision QA check on output slides")

    # slides build-rich-batch
    p = slides_sub.add_parser("build-rich-batch", help="Build multiple rich presentations in parallel")
    p.add_argument("--folder", help="Drive folder ID")
    p.add_argument("--link-share", action="store_true", help="Make link-shareable")
    p.add_argument("--parallel", type=int, default=4, help="Max parallel builds")
    p.add_argument("--spec", help="Inline JSON spec array")
    p.add_argument("--spec-file", help="Path to JSON spec file")

    # --- DOCS ---
    docs_p = sub.add_parser("docs", help="Google Docs operations")
    docs_sub = docs_p.add_subparsers(dest="docs_cmd")

    # docs create
    p = docs_sub.add_parser("create", help="Create a document")
    p.add_argument("title")

    # docs get
    p = docs_sub.add_parser("get", help="Get document structure")
    p.add_argument("document_id")

    # docs add-text
    p = docs_sub.add_parser("add-text", help="Add styled text")
    p.add_argument("document_id")
    p.add_argument("text")
    p.add_argument("--style", choices=["heading1", "heading2", "heading3", "heading4",
                                       "heading5", "heading6", "title", "subtitle", "normal"])
    p.add_argument("--text-style", help="JSON text style overrides")
    p.add_argument("--index", type=int, help="Insert at this document index (from docs get)")
    p.add_argument("--after-heading", help="Insert after heading containing this text")

    # docs add-table
    p = docs_sub.add_parser("add-table", help="Add a styled table")
    p.add_argument("document_id")
    p.add_argument("--data", required=True, help="2D JSON array")
    p.add_argument("--style", help="JSON style object")
    p.add_argument("--index", type=int, help="Insert at this document index")
    p.add_argument("--after-heading", help="Insert after heading containing this text")

    # docs add-image
    p = docs_sub.add_parser("add-image", help="Insert inline image")
    p.add_argument("document_id")
    p.add_argument("--url", required=True)
    p.add_argument("--width", type=float, help="Width in pt")
    p.add_argument("--height", type=float, help="Height in pt")
    p.add_argument("--index", type=int, help="Insert at this document index")
    p.add_argument("--after-heading", help="Insert after heading containing this text")

    # docs add-section
    p = docs_sub.add_parser("add-section", help="Add section break + heading")
    p.add_argument("document_id")
    p.add_argument("--title", required=True)

    # docs add-list
    p = docs_sub.add_parser("add-list", help="Add bulleted/numbered list")
    p.add_argument("document_id")
    p.add_argument("--items", required=True, help="JSON array of strings")
    p.add_argument("--type", dest="list_type", default="bullet", choices=["bullet", "numbered"])
    p.add_argument("--index", type=int, help="Insert at this document index")
    p.add_argument("--after-heading", help="Insert after heading containing this text")

    # docs add-page-break
    p = docs_sub.add_parser("add-page-break", help="Insert page break")
    p.add_argument("document_id")

    # docs add-divider
    p = docs_sub.add_parser("add-divider", help="Add visual divider")
    p.add_argument("document_id")

    # docs add-header
    p = docs_sub.add_parser("add-header", help="Add page header")
    p.add_argument("document_id")
    p.add_argument("--text", required=True)

    # docs add-footer
    p = docs_sub.add_parser("add-footer", help="Add page footer")
    p.add_argument("document_id")
    p.add_argument("--text", required=True)

    # docs replace-text
    p = docs_sub.add_parser("replace-text", help="Find and replace text")
    p.add_argument("document_id")
    p.add_argument("--find", required=True, help="Text to find")
    p.add_argument("--replace", required=True, help="Replacement text")
    p.add_argument("--ignore-case", action="store_true", help="Case-insensitive matching")

    # docs delete-range
    p = docs_sub.add_parser("delete-range", help="Delete content between indices")
    p.add_argument("document_id")
    p.add_argument("--start", type=int, required=True, dest="start_index", help="Start index")
    p.add_argument("--end", type=int, required=True, dest="end_index", help="End index")

    # docs export
    p = docs_sub.add_parser("export", help="Export document")
    p.add_argument("document_id")
    p.add_argument("--format", required=True, choices=["pdf", "docx", "md", "html", "epub"])
    p.add_argument("--output", help="Output file path")

    # docs build-rich
    p = docs_sub.add_parser("build-rich", help="Build rich document from JSON spec")
    p.add_argument("--title", help="Document title (overrides spec title)")
    p.add_argument("--spec", help="JSON spec as inline string")
    p.add_argument("--spec-file", dest="spec_file", help="JSON spec file path (or pipe via stdin)")
    p.add_argument("--folder", help="Drive folder ID")
    p.add_argument("--link-share", dest="link_share", action="store_true", help="Make link-shareable")
    p.add_argument("--share", help="Share with this email address")
    p.add_argument("--output", help="Save .docx locally (skip upload)")

    # docs build-rich-batch
    p = docs_sub.add_parser("build-rich-batch", help="Build multiple rich documents in parallel")
    p.add_argument("--specs", help="JSON array of specs as inline string")
    p.add_argument("--spec-file", dest="spec_file", help="JSON spec file path (or pipe via stdin)")
    p.add_argument("--folder", help="Drive folder ID")
    p.add_argument("--link-share", dest="link_share", action="store_true", help="Make link-shareable")
    p.add_argument("--share", help="Share with this email address")
    p.add_argument("--parallel", type=int, default=3, help="Max parallel workers (default: 3)")

    # --- SHEETS ---
    sheets_p = sub.add_parser("sheets", help="Google Sheets operations")
    sheets_sub = sheets_p.add_subparsers(dest="sheets_cmd")

    # sheets create
    p = sheets_sub.add_parser("create", help="Create a spreadsheet")
    p.add_argument("title")
    p.add_argument("--frozen-rows", type=int, help="Number of frozen header rows")

    # sheets get
    p = sheets_sub.add_parser("get", help="Get spreadsheet metadata")
    p.add_argument("spreadsheet_id")

    # sheets read
    p = sheets_sub.add_parser("read", help="Read spreadsheet data")
    p.add_argument("spreadsheet_id")
    p.add_argument("--range", help="A1 range notation")
    p.add_argument("--render", choices=["formatted", "raw", "formula"], default="formatted")

    # sheets write
    p = sheets_sub.add_parser("write", help="Write data to spreadsheet")
    p.add_argument("spreadsheet_id")
    p.add_argument("--range", required=True, help="Starting cell (e.g. A1)")
    p.add_argument("--data", required=True, help="2D JSON array")
    p.add_argument("--input", dest="input_type", choices=["raw", "user-entered"], default="user-entered")

    # sheets format
    p = sheets_sub.add_parser("format", help="Format cells")
    p.add_argument("spreadsheet_id")
    p.add_argument("--range", required=True)
    p.add_argument("--style", required=True, help="JSON style object")

    # sheets add-chart
    p = sheets_sub.add_parser("add-chart", help="Add native chart")
    p.add_argument("spreadsheet_id")
    p.add_argument("--range", required=True)
    p.add_argument("--type", dest="chart_type", default="bar",
                   choices=["bar", "line", "pie", "column", "area", "scatter"])

    # sheets add-tab
    p = sheets_sub.add_parser("add-tab", help="Add a new tab")
    p.add_argument("spreadsheet_id")
    p.add_argument("--title", dest="tab_title", required=True)
    p.add_argument("--color", help="Tab color (#RRGGBB)")

    # sheets delete-tab
    p = sheets_sub.add_parser("delete-tab", help="Delete a tab/sheet")
    p.add_argument("spreadsheet_id")
    p.add_argument("--sheet-id", type=int, help="Numeric sheet ID")
    p.add_argument("--title", help="Tab title (alternative to --sheet-id)")

    # sheets conditional-format
    p = sheets_sub.add_parser("conditional-format", help="Add conditional formatting")
    p.add_argument("spreadsheet_id")
    p.add_argument("--range", required=True)
    p.add_argument("--rule", required=True, help="JSON rule object")

    # sheets auto-resize
    p = sheets_sub.add_parser("auto-resize", help="Auto-resize columns")
    p.add_argument("spreadsheet_id")
    p.add_argument("--columns", help="Column range (e.g. 0:10)")

    # sheets build-rich
    p = sheets_sub.add_parser("build-rich", help="Build rich spreadsheet from JSON spec")
    p.add_argument("--title", help="Spreadsheet title (overrides spec title)")
    p.add_argument("--folder", help="Drive folder ID")
    p.add_argument("--link-share", dest="link_share", action="store_true", help="Make link-shareable")
    p.add_argument("--share", help="Share with this email address")
    p.add_argument("--output", help="Save .xlsx locally instead of only uploading")
    p.add_argument("--spec-file", dest="spec_file", help="JSON spec file path (or pipe via stdin)")

    # sheets build-rich-batch
    p = sheets_sub.add_parser("build-rich-batch", help="Build multiple rich spreadsheets in parallel")
    p.add_argument("--folder", help="Drive folder ID")
    p.add_argument("--link-share", dest="link_share", action="store_true", help="Make link-shareable")
    p.add_argument("--parallel", type=int, default=4, help="Max parallel workers (default: 4)")
    p.add_argument("--spec-file", dest="spec_file", help="JSON spec file path (or pipe via stdin)")

    # sheets analyze
    p = sheets_sub.add_parser("analyze", help="Analyze spreadsheet structure and data")
    p.add_argument("spreadsheet_id", nargs="?", help="Google Sheets ID")
    p.add_argument("--file", help="Local .xlsx file path")
    p.add_argument("--sheet", help="Specific sheet name to analyze")

    # sheets query
    p = sheets_sub.add_parser("query", help="Filter/aggregate spreadsheet data")
    p.add_argument("spreadsheet_id", nargs="?", help="Google Sheets ID")
    p.add_argument("--file", help="Local .xlsx file path")
    p.add_argument("--filter", help="JSON filter conditions")
    p.add_argument("--groupby", help="Column name to group by")
    p.add_argument("--agg", help="JSON aggregation spec (e.g. '{\"Revenue\": \"sum\"}')")
    p.add_argument("--sort", help="Column to sort by (prefix - for descending)")
    p.add_argument("--limit", type=int, help="Max rows to return")
    p.add_argument("--sheet", help="Specific sheet name")

    # sheets download
    p = sheets_sub.add_parser("download", help="Download Google Sheet as .xlsx")
    p.add_argument("spreadsheet_id")
    p.add_argument("--output", help="Local output path")

    # sheets edit
    p = sheets_sub.add_parser("edit", help="Edit existing spreadsheet (download, modify, re-upload)")
    p.add_argument("spreadsheet_id")
    p.add_argument("--spec-file", dest="spec_file", help="JSON edit spec file (or pipe via stdin)")
    p.add_argument("--link-share", dest="link_share", action="store_true", help="Make link-shareable")

    # --- DRIVE ---
    drive_p = sub.add_parser("drive", help="Google Drive operations")
    drive_sub = drive_p.add_subparsers(dest="drive_cmd")

    # drive list
    p = drive_sub.add_parser("list", help="List files")
    p.add_argument("--folder", help="Parent folder ID")
    p.add_argument("--type", dest="file_type", choices=["slides", "docs", "sheets"])

    # drive share
    p = drive_sub.add_parser("share", help="Share a file")
    p.add_argument("file_id")
    p.add_argument("--email", required=True)
    p.add_argument("--role", default="reader", choices=["reader", "writer", "commenter"])

    # drive move
    p = drive_sub.add_parser("move", help="Move a file")
    p.add_argument("file_id")
    p.add_argument("--to", required=True, help="Destination folder ID")

    # drive upload
    p = drive_sub.add_parser("upload", help="Upload a local file")
    p.add_argument("local_path")
    p.add_argument("--folder", help="Destination folder ID")
    p.add_argument("--mime-type", help="MIME type override")

    # drive create-folder
    p = drive_sub.add_parser("create-folder", help="Create a folder")
    p.add_argument("name")
    p.add_argument("--parent", help="Parent folder ID")

    # drive rename
    p = drive_sub.add_parser("rename", help="Rename a file or folder")
    p.add_argument("file_id")
    p.add_argument("new_name")

    # drive search
    p = drive_sub.add_parser("search", help="Search files by name")
    p.add_argument("query", help="Search term (matches name contains)")
    p.add_argument("--parent", help="Filter to parent folder ID")
    p.add_argument("--folder-only", action="store_true", help="Only return folders")
    p.add_argument("--limit", type=int, help="Max results (default: 25)")

    return parser


# ============================================================
# Dispatch
# ============================================================

DISPATCH = {
    ("auth", "status"): cmd_auth_status,
    ("slides", "create"): cmd_slides_create,
    ("slides", "get"): cmd_slides_get,
    ("slides", "add-slide"): cmd_slides_add_slide,
    ("slides", "add-text"): cmd_slides_add_text,
    ("slides", "add-table"): cmd_slides_add_table,
    ("slides", "add-image"): cmd_slides_add_image,
    ("slides", "add-shape"): cmd_slides_add_shape,
    ("slides", "set-background"): cmd_slides_set_background,
    ("slides", "apply-theme"): cmd_slides_apply_theme,
    ("slides", "add-chart"): cmd_slides_add_chart,
    ("slides", "add-layout"): cmd_slides_add_layout,
    ("slides", "export"): cmd_slides_export,
    ("slides", "batch-update"): cmd_slides_batch_update,
    ("slides", "build-rich"): cmd_slides_build_rich,
    ("slides", "build-rich-batch"): cmd_slides_build_rich_batch,
    ("docs", "create"): cmd_docs_create,
    ("docs", "get"): cmd_docs_get,
    ("docs", "add-text"): cmd_docs_add_text,
    ("docs", "add-table"): cmd_docs_add_table,
    ("docs", "add-image"): cmd_docs_add_image,
    ("docs", "add-section"): cmd_docs_add_section,
    ("docs", "add-list"): cmd_docs_add_list,
    ("docs", "add-page-break"): cmd_docs_add_page_break,
    ("docs", "add-divider"): cmd_docs_add_divider,
    ("docs", "add-header"): cmd_docs_add_header,
    ("docs", "add-footer"): cmd_docs_add_footer,
    ("docs", "replace-text"): cmd_docs_replace_text,
    ("docs", "delete-range"): cmd_docs_delete_range,
    ("docs", "export"): cmd_docs_export,
    ("docs", "build-rich"): cmd_docs_build_rich,
    ("docs", "build-rich-batch"): cmd_docs_build_rich_batch,
    ("sheets", "create"): cmd_sheets_create,
    ("sheets", "get"): cmd_sheets_get,
    ("sheets", "read"): cmd_sheets_read,
    ("sheets", "write"): cmd_sheets_write,
    ("sheets", "format"): cmd_sheets_format,
    ("sheets", "add-chart"): cmd_sheets_add_chart,
    ("sheets", "add-tab"): cmd_sheets_add_tab,
    ("sheets", "delete-tab"): cmd_sheets_delete_tab,
    ("sheets", "conditional-format"): cmd_sheets_conditional_format,
    ("sheets", "auto-resize"): cmd_sheets_auto_resize,
    ("sheets", "build-rich"): cmd_sheets_build_rich,
    ("sheets", "build-rich-batch"): cmd_sheets_build_rich_batch,
    ("sheets", "analyze"): cmd_sheets_analyze,
    ("sheets", "query"): cmd_sheets_query,
    ("sheets", "download"): cmd_sheets_download,
    ("sheets", "edit"): cmd_sheets_edit,
    ("drive", "list"): cmd_drive_list,
    ("drive", "share"): cmd_drive_share,
    ("drive", "move"): cmd_drive_move,
    ("drive", "upload"): cmd_drive_upload,
    ("drive", "create-folder"): cmd_drive_create_folder,
    ("drive", "rename"): cmd_drive_rename,
    ("drive", "search"): cmd_drive_search,
}


def main():
    parser = build_parser()
    args = parser.parse_args()

    if not args.service:
        parser.print_help()
        sys.exit(1)

    # Determine subcommand
    cmd_attr = f"{args.service}_cmd"
    sub_cmd = getattr(args, cmd_attr, None)
    if not sub_cmd:
        parser.parse_args([args.service, "--help"])
        sys.exit(1)

    key = (args.service, sub_cmd)
    handler = DISPATCH.get(key)
    if not handler:
        print(f"ERROR: Unknown command: {args.service} {sub_cmd}", file=sys.stderr)
        sys.exit(1)

    handler(args)


if __name__ == "__main__":
    main()
